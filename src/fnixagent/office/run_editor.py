"""Run 级原地修改器(Phase 5.3)。

对 Word/Excel/PPT 文档执行 Run 级原地编辑操作,保留原有 run 边界与格式。
支持四种操作类型:replace/insert/delete/fill_blank。

参考实现:
  - Office-Word-MCP-Server/word_document_server/utils/document_utils.py
    find_and_replace_text(逐 run 替换,保留格式)
  - fnixagent.office.word.WordExpert._delete_paragraph
    para._element.getparent().remove(para._element)

设计要点:
  - 替换时遍历 para.runs,在每个 run 的 text 上做 replace,
    避免跨 run 拼接丢失格式
  - 填括号:用正则定位 (\\s*) 或 (\\s*),在对应 run 上替换
  - 删除段落:用 para._element.getparent().remove(para._element)
  - 所有操作失败不抛异常,记录到 failed_ops,继续后续操作
"""
from __future__ import annotations

import os
import re
from dataclasses import asdict, dataclass, field
from typing import Any, Optional

from fnixagent.office.base import BaseExpert, ExpertError, ExpertResult


# 全角/半角空白括号正则:匹配 () 或 (  ) 或 （） 或 （  ）
_PAREN_BLANK_RE = re.compile(r"[（(]\s*[）)]")


# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------


@dataclass
class EditOp:
    """单个编辑操作。

    Attributes:
        op_type: 操作类型("replace"/"insert"/"delete"/"fill_blank")
        target: 定位文本/段落索引/括号位置/单元格引用
        value: 替换/插入的值(delete 时可为空)
        position: insert 位置("before"/"after"/"replace")
        preserve_format: 是否保留原格式
    """

    op_type: str  # "replace"/"insert"/"delete"/"fill_blank"
    target: str  # 定位文本/段落索引/括号位置
    value: str = ""  # 替换/插入的值
    position: str = "after"  # insert 位置: before/after/replace
    preserve_format: bool = True  # 是否保留原格式


@dataclass
class EditReport:
    """编辑结果报告。

    Attributes:
        total_ops: 操作总数
        applied_ops: 成功执行的操作数
        failed_ops: 失败的操作数
        details: 每个操作的执行记录 dict 列表
    """

    total_ops: int = 0
    applied_ops: int = 0
    failed_ops: int = 0
    details: list[dict] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Run 级原地修改器
# ---------------------------------------------------------------------------


class RunEditor(BaseExpert):
    """Run 级原地修改器(Phase 5.3)。

    对 Word/Excel/PPT 文档执行 Run 级原地编辑,保留原有 run 边界与格式。
    支持 replace/insert/delete/fill_blank 四种操作。

    能力边界:
      - 仅处理 .docx/.xlsx/.pptx
      - 跨 run 文本不拼接(逐 run 替换,可能漏匹配跨 run 关键词)
      - Excel insert 视作单元格赋值
      - PPT 段落删除为清空 run 文本(保留 shape)
    """

    @property
    def name(self) -> str:
        return "run_editor"

    # ------------------------------------------------------------------
    # 公共 API
    # ------------------------------------------------------------------

    def edit(
        self,
        path: str,
        ops: list[EditOp],
        output_path: Optional[str] = None,
    ) -> ExpertResult:
        """按扩展名派发到对应文档类型的编辑方法。

        Args:
            path: 输入文档路径
            ops: 编辑操作列表
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": {...}})
        """
        ext = os.path.splitext(path)[1].lstrip(".").lower()
        if ext == "docx":
            return self.edit_word(path, ops, output_path)
        if ext == "xlsx":
            return self.edit_excel(path, ops, output_path)
        if ext == "pptx":
            return self.edit_ppt(path, ops, output_path)
        return self._failure(
            f"unsupported extension: .{ext}, allowed: docx/xlsx/pptx"
        )

    def edit_word(
        self,
        path: str,
        ops: list[EditOp],
        output_path: Optional[str] = None,
    ) -> ExpertResult:
        """Word run 级编辑。

        Args:
            path: 输入 .docx 路径
            ops: 编辑操作列表,每项为 EditOp
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": EditReport dict})
        """
        if not isinstance(ops, list) or not ops:
            return self._failure("ops must be a non-empty list")

        err = self._validate_path(
            path, must_exist=True, allowed_exts=("docx",)
        )
        if err:
            return self._failure(err)
        if output_path:
            err = self._validate_path(output_path, allowed_exts=("docx",))
            if err:
                return self._failure(err)

        try:
            self._require_lib("docx")
            from docx import Document
        except ExpertError as e:
            return self._failure(str(e))

        try:
            doc = Document(path)
            report = EditReport(total_ops=len(ops))

            # 收集所有可编辑段落(正文 + 表格内)
            all_paragraphs = list(doc.paragraphs)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        all_paragraphs.extend(cell.paragraphs)

            for op in ops:
                self._apply_word_op(doc, all_paragraphs, op, report)

            save_path = output_path or path
            doc.save(save_path)
            return self._success(save_path, report=asdict(report))
        except (PermissionError, IOError) as e:
            return self._failure(f"edit_word IO failed: {e}")
        except Exception as e:
            return self._failure(f"edit_word failed: {e}")

    def edit_excel(
        self,
        path: str,
        ops: list[EditOp],
        output_path: Optional[str] = None,
    ) -> ExpertResult:
        """Excel 单元格编辑。

        - target 为单元格引用(如 "A1")时:直接定位
        - target 为普通文本时:在所有 sheet 中查找替换

        Args:
            path: 输入 .xlsx 路径
            ops: 编辑操作列表
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": EditReport dict})
        """
        if not isinstance(ops, list) or not ops:
            return self._failure("ops must be a non-empty list")

        err = self._validate_path(
            path, must_exist=True, allowed_exts=("xlsx",)
        )
        if err:
            return self._failure(err)
        if output_path:
            err = self._validate_path(output_path, allowed_exts=("xlsx",))
            if err:
                return self._failure(err)

        try:
            self._require_lib("openpyxl")
            from openpyxl import load_workbook
        except ExpertError as e:
            return self._failure(str(e))

        try:
            wb = load_workbook(path)
            report = EditReport(total_ops=len(ops))

            for op in ops:
                self._apply_excel_op(wb, op, report)

            save_path = output_path or path
            wb.save(save_path)
            return self._success(save_path, report=asdict(report))
        except (PermissionError, IOError) as e:
            return self._failure(f"edit_excel IO failed: {e}")
        except Exception as e:
            return self._failure(f"edit_excel failed: {e}")

    def edit_ppt(
        self,
        path: str,
        ops: list[EditOp],
        output_path: Optional[str] = None,
    ) -> ExpertResult:
        """PPT 文本框编辑。

        Args:
            path: 输入 .pptx 路径
            ops: 编辑操作列表
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": EditReport dict})
        """
        if not isinstance(ops, list) or not ops:
            return self._failure("ops must be a non-empty list")

        err = self._validate_path(
            path, must_exist=True, allowed_exts=("pptx",)
        )
        if err:
            return self._failure(err)
        if output_path:
            err = self._validate_path(output_path, allowed_exts=("pptx",))
            if err:
                return self._failure(err)

        try:
            self._require_lib("pptx")
            from pptx import Presentation
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation(path)
            report = EditReport(total_ops=len(ops))

            for op in ops:
                self._apply_ppt_op(prs, op, report)

            save_path = output_path or path
            prs.save(save_path)
            return self._success(save_path, report=asdict(report))
        except (PermissionError, IOError) as e:
            return self._failure(f"edit_ppt IO failed: {e}")
        except Exception as e:
            return self._failure(f"edit_ppt failed: {e}")

    # ------------------------------------------------------------------
    # Word 内部方法
    # ------------------------------------------------------------------

    def _apply_word_op(
        self,
        doc: Any,
        paragraphs: list,
        op: EditOp,
        report: EditReport,
    ) -> None:
        """对 Word 文档应用单个编辑操作,结果累计到 report。"""
        try:
            if op.op_type == "replace":
                count = 0
                for para in paragraphs:
                    if self._word_replace_in_paragraph(
                        para, op.target, op.value, op.preserve_format
                    ):
                        count += 1
                self._record_op(report, op, count > 0, count=count)

            elif op.op_type == "insert":
                target_para = None
                for para in paragraphs:
                    if op.target in (para.text or ""):
                        target_para = para
                        break
                if target_para is None:
                    self._record_op(
                        report, op, False, error="target paragraph not found"
                    )
                    return

                success = False
                if op.position == "replace":
                    # 替换段落文本:清空所有 run,再添加一个新 run
                    for run in list(target_para.runs):
                        run._element.getparent().remove(run._element)
                    target_para.add_run(op.value)
                    success = True
                elif op.position == "before":
                    new_para = doc.add_paragraph(op.value)
                    target_para._element.addprevious(new_para._element)
                    success = True
                else:  # after
                    success = self._word_insert_after(
                        doc, target_para, op.value
                    )
                self._record_op(report, op, success, position=op.position)

            elif op.op_type == "delete":
                deleted = 0
                # 用 list() 拷贝,避免遍历时修改
                for para in list(paragraphs):
                    if op.target in (para.text or ""):
                        if self._word_delete_paragraph(doc, para):
                            deleted += 1
                            paragraphs.remove(para)
                self._record_op(report, op, deleted > 0, count=deleted)

            elif op.op_type == "fill_blank":
                filled = 0
                for para in paragraphs:
                    if self._word_fill_paren(para, op.target, op.value):
                        filled += 1
                self._record_op(report, op, filled > 0, count=filled)

            else:
                self._record_op(
                    report,
                    op,
                    False,
                    error=f"unknown op_type: {op.op_type}",
                )
        except Exception as e:
            self._record_op(report, op, False, error=str(e))

    def _word_replace_in_paragraph(
        self,
        para: Any,
        old: str,
        new: str,
        preserve_format: bool = True,
    ) -> bool:
        """段落内替换,逐 run 操作保留格式。

        参考: Office-Word-MCP-Server/word_document_server/utils/document_utils.py
        的 find_and_replace_text 实现。

        Args:
            para: python-docx Paragraph 对象
            old: 待查找文本
            new: 替换为的文本
            preserve_format: 是否保留原格式(逐 run 替换天然保留)

        Returns:
            是否发生替换
        """
        if not old:
            return False
        if old not in (para.text or ""):
            return False

        replaced = False
        # 逐 run 替换:在每个 run 的 text 上做 replace,保留 run 边界
        # 局限:不处理跨 run 的关键词
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                replaced = True
        return replaced

    def _word_fill_paren(
        self, para: Any, blank_marker: str, value: str
    ) -> bool:
        """填入括号空白,如把"（ ）"变成"（B）"。

        用正则定位 全角/半角 空白括号,在对应 run 上替换。

        Args:
            para: python-docx Paragraph 对象
            blank_marker: 括号标记(如 "（）"/"()"),用于判断全角/半角
            value: 填入的值

        Returns:
            是否填入成功
        """
        # 判断目标使用全角还是半角括号
        use_fullwidth = "（" in blank_marker or "）" in blank_marker
        replacement = f"（{value}）" if use_fullwidth else f"({value})"

        filled = False
        for run in para.runs:
            if _PAREN_BLANK_RE.search(run.text):
                run.text = _PAREN_BLANK_RE.sub(
                    replacement, run.text, count=1
                )
                filled = True
                break
        return filled

    def _word_delete_paragraph(self, doc: Any, para: Any) -> bool:
        """安全删除段落。

        参考: fnixagent.office.word.WordExpert._delete_paragraph
        用 para._element.getparent().remove(para._element)

        Args:
            doc: python-docx Document 对象(保留接口,暂不使用)
            para: 待删除的 Paragraph 对象

        Returns:
            是否删除成功
        """
        try:
            p = para._element
            parent = p.getparent()
            if parent is None:
                return False
            parent.remove(p)
            return True
        except Exception:
            return False

    def _word_insert_after(
        self,
        doc: Any,
        target_para: Any,
        text: str,
        style: Optional[str] = None,
    ) -> bool:
        """在目标段落后插入新段落。

        Args:
            doc: python-docx Document 对象
            target_para: 锚点段落
            text: 新段落文本
            style: 段落样式名;None 使用默认

        Returns:
            是否插入成功
        """
        try:
            if style:
                new_para = doc.add_paragraph(text, style=style)
            else:
                new_para = doc.add_paragraph(text)
            target_para._element.addnext(new_para._element)
            return True
        except Exception:
            return False

    # ------------------------------------------------------------------
    # Excel 内部方法
    # ------------------------------------------------------------------

    def _apply_excel_op(
        self,
        wb: Any,
        op: EditOp,
        report: EditReport,
    ) -> None:
        """对 Excel 工作簿应用单个编辑操作,结果累计到 report。"""
        try:
            if op.op_type == "replace":
                if self._is_cell_ref(op.target):
                    ws = wb.active
                    success = self._excel_set_cell(
                        ws, op.target, op.value, op.preserve_format
                    )
                    self._record_op(report, op, success)
                else:
                    # 在所有 sheet 中查找替换
                    count = 0
                    for ws in wb.worksheets:
                        for row in ws.iter_rows():
                            for cell in row:
                                if (
                                    cell.value is not None
                                    and isinstance(cell.value, str)
                                    and op.target in cell.value
                                ):
                                    cell.value = cell.value.replace(
                                        op.target, op.value
                                    )
                                    count += 1
                    self._record_op(report, op, count > 0, count=count)

            elif op.op_type == "insert":
                ws = wb.active
                if self._is_cell_ref(op.target):
                    success = self._excel_set_cell(
                        ws, op.target, op.value, op.preserve_format
                    )
                    self._record_op(report, op, success)
                else:
                    self._record_op(
                        report,
                        op,
                        False,
                        error="insert requires a cell reference as target",
                    )

            elif op.op_type == "delete":
                ws = wb.active
                if self._is_cell_ref(op.target):
                    try:
                        ws[op.target] = None
                        self._record_op(report, op, True)
                    except Exception as e:
                        self._record_op(report, op, False, error=str(e))
                else:
                    self._record_op(
                        report,
                        op,
                        False,
                        error="delete requires a cell reference as target",
                    )

            elif op.op_type == "fill_blank":
                ws = wb.active
                if self._is_cell_ref(op.target):
                    success = self._excel_set_cell(
                        ws, op.target, op.value, op.preserve_format
                    )
                    self._record_op(report, op, success)
                else:
                    self._record_op(
                        report,
                        op,
                        False,
                        error="fill_blank requires a cell reference as target",
                    )

            else:
                self._record_op(
                    report,
                    op,
                    False,
                    error=f"unknown op_type: {op.op_type}",
                )
        except Exception as e:
            self._record_op(report, op, False, error=str(e))

    def _excel_set_cell(
        self,
        ws: Any,
        cell_ref: str,
        value: Any,
        preserve_format: bool = True,
    ) -> bool:
        """设置单元格值保留格式。

        Args:
            ws: openpyxl Worksheet 对象
            cell_ref: 单元格引用(如 "A1")
            value: 新值
            preserve_format: 是否保留原格式(对齐/边框/填充/数字格式)

        Returns:
            是否设置成功
        """
        try:
            cell = ws[cell_ref]
            # 保留原格式
            old_font = cell.font
            old_fill = cell.fill
            old_align = cell.alignment
            old_border = cell.border
            old_num_fmt = cell.number_format

            cell.value = value

            if preserve_format:
                # openpyxl 中 Font/Fill 等是不可变对象,赋值即可
                cell.font = old_font
                cell.fill = old_fill
                cell.alignment = old_align
                cell.border = old_border
                cell.number_format = old_num_fmt
            return True
        except Exception:
            return False

    # ------------------------------------------------------------------
    # PPT 内部方法
    # ------------------------------------------------------------------

    def _apply_ppt_op(
        self,
        prs: Any,
        op: EditOp,
        report: EditReport,
    ) -> None:
        """对 PPT 文稿应用单个编辑操作,结果累计到 report。"""
        try:
            if op.op_type == "replace":
                count = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if op.target in (run.text or ""):
                                    run.text = run.text.replace(
                                        op.target, op.value
                                    )
                                    count += 1
                self._record_op(report, op, count > 0, count=count)

            elif op.op_type == "insert":
                success = False
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            tf = shape.text_frame
                            p = tf.add_paragraph()
                            p.text = op.value
                            success = True
                            break
                    if success:
                        break
                self._record_op(report, op, success)

            elif op.op_type == "delete":
                count = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in list(shape.text_frame.paragraphs):
                            if op.target in (para.text or ""):
                                # 清空该段落的 run 文本
                                # (python-pptx 不支持直接删除段落,清空替代)
                                for run in para.runs:
                                    run.text = ""
                                count += 1
                self._record_op(report, op, count > 0, count=count)

            elif op.op_type == "fill_blank":
                count = 0
                use_fullwidth = "（" in op.target or "）" in op.target
                replacement = (
                    f"（{op.value}）" if use_fullwidth else f"({op.value})"
                )
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if _PAREN_BLANK_RE.search(run.text or ""):
                                    run.text = _PAREN_BLANK_RE.sub(
                                        replacement, run.text, count=1
                                    )
                                    count += 1
                self._record_op(report, op, count > 0, count=count)

            else:
                self._record_op(
                    report,
                    op,
                    False,
                    error=f"unknown op_type: {op.op_type}",
                )
        except Exception as e:
            self._record_op(report, op, False, error=str(e))

    # ------------------------------------------------------------------
    # 通用工具
    # ------------------------------------------------------------------

    @staticmethod
    def _is_cell_ref(s: str) -> bool:
        """判断字符串是否为 Excel 单元格引用(如 A1, B2, AA10)。

        Args:
            s: 待判断字符串

        Returns:
            True 视为单元格引用
        """
        if not s or not isinstance(s, str):
            return False
        return bool(re.match(r"^[A-Za-z]+[0-9]+$", s.strip()))

    @staticmethod
    def _record_op(
        report: EditReport,
        op: EditOp,
        success: bool,
        **extra: Any,
    ) -> None:
        """记录单个操作结果到 report。

        Args:
            report: EditReport 实例
            op: 原始 EditOp
            success: 是否成功
            **extra: 附加记录字段(count/position/error 等)
        """
        if success:
            report.applied_ops += 1
        else:
            report.failed_ops += 1
        detail: dict = {
            "op_type": op.op_type,
            "target": op.target,
            "success": success,
        }
        if op.value:
            detail["value"] = op.value
        detail.update(extra)
        report.details.append(detail)
