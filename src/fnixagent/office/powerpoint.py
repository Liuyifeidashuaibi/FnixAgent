"""PPT Expert(P2-9)。

PowerPoint 演示文稿创建/slide/主题/图片/图表/导出图片。

专家职责:
  - 创建/编辑 .pptx(slide/标题/项目符号/主题/图片/图表)
  - 导出每页为图片(Windows COM 或 LibreOffice 兜底)

底层依赖:
  - python-pptx(可选,创建/编辑)
  - pywin32(可选,Windows COM 导出图片)
  - PyMuPDF/fitz(可选,LibreOffice 路径下 PDF 转图片)

降级策略:
  - 依赖缺失 → ExpertError 提示安装
  - COM 调用后必须 Quit + CoUninitialize 释放
  - LibreOffice headless 失败时返回 _failure,不崩溃
"""
from __future__ import annotations

import os
from typing import Any, Optional

from fnixagent.office.base import BaseExpert, ExpertError, ExpertResult


class PPTExpert(BaseExpert):
    """PowerPoint 演示文稿专家。

    全部方法返回 ExpertResult。

    能力边界:
      - 仅处理 .pptx(python-pptx 不支持 .ppt 旧格式)
      - 导出图片依赖 PowerPoint(COM)或 LibreOffice,纯 Python 无直接方案
      - 主题应用为简化实现(全局颜色/字体覆盖)
    """

    @property
    def name(self) -> str:
        return "ppt"

    # ------------------------------------------------------------------
    # 创建
    # ------------------------------------------------------------------

    def create(
        self,
        output_path: str,
        title: str = "",
        subtitle: str = "",
        slides: Optional[list[dict]] = None,
    ) -> ExpertResult:
        """创建 PPT 文稿。

        Args:
            output_path: 输出 .pptx 路径
            title: 首页标题
            subtitle: 首页副标题
            slides: 后续页内容,每项 {layout, title, content, bullets}

        Returns:
            ExpertResult(output=output_path, slides=N)
        """
        err = self._validate_path(output_path, allowed_exts=("pptx",))
        if err:
            return self._failure(err)

        try:
            self._require_lib("pptx")
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation()
            # 首页(标题页)
            if title:
                slide_layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = title
                if len(slide.placeholders) > 1 and subtitle:
                    slide.placeholders[1].text = subtitle

            # 后续页
            if slides:
                for spec in slides:
                    layout_idx = spec.get("layout", 1)  # 默认 Title and Content
                    try:
                        slide_layout = prs.slide_layouts[layout_idx]
                    except IndexError:
                        slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    s_title = spec.get("title")
                    if s_title and slide.shapes.title:
                        slide.shapes.title.text = s_title
                    bullets = spec.get("bullets") or []
                    content_text = spec.get("content", "")
                    # 优先写入 placeholder
                    if slide.placeholders:
                        for ph in slide.placeholders:
                            if ph.placeholder_format.idx == 1:  # Content placeholder
                                if bullets:
                                    tf = ph.text_frame
                                    tf.text = bullets[0] if bullets else ""
                                    for b in bullets[1:]:
                                        p = tf.add_paragraph()
                                        p.text = b
                                elif content_text:
                                    ph.text = content_text
                                break

            prs.save(output_path)
            return self._success(output_path, slides=len(prs.slides))
        except (PermissionError, IOError) as e:
            return self._failure(f"create ppt IO failed: {e}")
        except Exception as e:
            return self._failure(f"create ppt failed: {e}")

    # ------------------------------------------------------------------
    # 添加 slide
    # ------------------------------------------------------------------

    def add_slide(
        self,
        path: str,
        title: str = "",
        bullets: Optional[list[str]] = None,
        layout: int = 1,
    ) -> ExpertResult:
        """向现有 PPT 追加一页。

        Args:
            path: .pptx 路径(原地修改)
            title: 标题
            bullets: 项目符号内容列表
            layout: slide layout 索引(默认 1=Title and Content)

        Returns:
            ExpertResult(output=slide_index)
        """
        err = self._validate_path(
            path, must_exist=True, allowed_exts=("pptx",)
        )
        if err:
            return self._failure(err)
        err = self._validate_int(layout, "layout", min_value=0)
        if err:
            return self._failure(err)

        try:
            self._require_lib("pptx")
            from pptx import Presentation
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation(path)
            try:
                slide_layout = prs.slide_layouts[layout]
            except IndexError:
                slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            if title and slide.shapes.title:
                slide.shapes.title.text = title
            if bullets:
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        tf = ph.text_frame
                        tf.text = bullets[0]
                        for b in bullets[1:]:
                            p = tf.add_paragraph()
                            p.text = b
                        break
            prs.save(path)
            return self._success(len(prs.slides) - 1)
        except (PermissionError, IOError) as e:
            return self._failure(f"add_slide IO failed: {e}")
        except Exception as e:
            return self._failure(f"add_slide failed: {e}")

    # ------------------------------------------------------------------
    # 主题
    # ------------------------------------------------------------------

    def apply_theme(
        self,
        path: str,
        theme_color: str = "1F4E79",
        font_name: str = "",
        font_size_title: int = 32,
        font_size_body: int = 18,
    ) -> ExpertResult:
        """应用基础主题(颜色 + 字体)。

        Args:
            path: .pptx 路径(原地修改)
            theme_color: 主色(HEX,无 #)
            font_name: 字体名(空串保留默认)
            font_size_title: 标题字号
            font_size_body: 正文字号

        Returns:
            ExpertResult(output=theme_color)
        """
        err = self._validate_path(
            path, must_exist=True, allowed_exts=("pptx",)
        )
        if err:
            return self._failure(err)
        err = self._validate_string(theme_color, "theme_color")
        if err:
            return self._failure(err)
        err = self._validate_int(font_size_title, "font_size_title", min_value=1)
        if err:
            return self._failure(err)
        err = self._validate_int(font_size_body, "font_size_body", min_value=1)
        if err:
            return self._failure(err)

        try:
            self._require_lib("pptx")
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation(path)
            color = RGBColor.from_string(theme_color)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = color
                            if font_name:
                                run.font.name = font_name
                            # 简化判断:首段视为标题
                            if para == shape.text_frame.paragraphs[0] and shape == slide.shapes[0]:
                                run.font.size = Pt(font_size_title)
                            else:
                                run.font.size = Pt(font_size_body)
            prs.save(path)
            return self._success(theme_color, font=font_name)
        except (PermissionError, IOError) as e:
            return self._failure(f"apply_theme IO failed: {e}")
        except Exception as e:
            return self._failure(f"apply_theme failed: {e}")

    # ------------------------------------------------------------------
    # 图片
    # ------------------------------------------------------------------

    def insert_image(
        self,
        path: str,
        image_path: str,
        slide_index: int,
        left: float = 1.0,
        top: float = 2.0,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> ExpertResult:
        """向指定 slide 插入图片。

        Args:
            path: .pptx 路径(原地修改)
            image_path: 图片路径
            slide_index: 目标 slide 索引(0-based)
            left/top: 位置(英寸)
            width/height: 尺寸(英寸,None 保持原比例)

        Returns:
            ExpertResult(output=slide_index)
        """
        err = self._validate_path(
            path, must_exist=True, allowed_exts=("pptx",)
        )
        if err:
            return self._failure(err)
        # 图片格式校验:常见光栅/矢量图
        err = self._validate_path(
            image_path,
            must_exist=True,
            allowed_exts=("png", "jpg", "jpeg", "gif", "bmp", "svg", "tif", "tiff"),
        )
        if err:
            return self._failure(err)
        err = self._validate_int(slide_index, "slide_index", min_value=0)
        if err:
            return self._failure(err)

        try:
            self._require_lib("pptx")
            from pptx import Presentation
            from pptx.util import Inches
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation(path)
            if slide_index < 0 or slide_index >= len(prs.slides):
                return self._failure(f"slide_index out of range: {slide_index}")
            slide = prs.slides[slide_index]
            kwargs = {"left": Inches(left), "top": Inches(top)}
            if width is not None:
                kwargs["width"] = Inches(width)
            if height is not None:
                kwargs["height"] = Inches(height)
            slide.shapes.add_picture(image_path, **kwargs)
            prs.save(path)
            return self._success(slide_index)
        except (PermissionError, IOError) as e:
            return self._failure(f"insert_image IO failed: {e}")
        except Exception as e:
            return self._failure(f"insert_image failed: {e}")

    # ------------------------------------------------------------------
    # 图表
    # ------------------------------------------------------------------

    def insert_chart(
        self,
        path: str,
        slide_index: int,
        chart_type: str,
        categories: list[str],
        series: dict[str, list[float]],
        title: str = "",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 6.0,
        height: float = 4.0,
    ) -> ExpertResult:
        """向指定 slide 插入图表。

        Args:
            path: .pptx 路径(原地修改)
            slide_index: 目标 slide 索引(0-based)
            chart_type: bar / line / pie
            categories: 类别标签
            series: {series_name: [values]}
            title: 图表标题
            left/top/width/height: 位置与尺寸(英寸)

        Returns:
            ExpertResult(output=slide_index)
        """
        err = self._validate_path(
            path, must_exist=True, allowed_exts=("pptx",)
        )
        if err:
            return self._failure(err)
        err = self._validate_int(slide_index, "slide_index", min_value=0)
        if err:
            return self._failure(err)
        err = self._validate_string(chart_type, "chart_type")
        if err:
            return self._failure(err)

        try:
            self._require_lib("pptx")
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation(path)
            if slide_index < 0 or slide_index >= len(prs.slides):
                return self._failure(f"slide_index out of range: {slide_index}")
            slide = prs.slides[slide_index]

            chart_data = CategoryChartData()
            chart_data.categories = categories
            for sname, values in series.items():
                chart_data.add_series(sname, values)

            ct = chart_type.lower()
            chart_type_map = {
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
            }
            if ct not in chart_type_map:
                return self._failure(f"unsupported chart_type: {chart_type}")
            chart_frame = slide.shapes.add_chart(
                chart_type_map[ct],
                Inches(left), Inches(top), Inches(width), Inches(height),
                chart_data,
            )
            chart = chart_frame.chart
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
            else:
                chart.has_title = False
            prs.save(path)
            return self._success(slide_index, chart_type=ct)
        except (PermissionError, IOError) as e:
            return self._failure(f"insert_chart IO failed: {e}")
        except Exception as e:
            return self._failure(f"insert_chart failed: {e}")

    # ------------------------------------------------------------------
    # 导出图片
    # ------------------------------------------------------------------

    def export_images(
        self,
        path: str,
        output_dir: str,
        slide_indices: Optional[list[int]] = None,
    ) -> ExpertResult:
        """将 PPT 每页导出为图片(需 LibreOffice 或 COM 接口)。

        Args:
            path: .pptx 路径
            output_dir: 输出目录
            slide_indices: 指定页索引;None 全部

        Returns:
            ExpertResult(output=[image_paths])

        Note:
            纯 Python 无直接方案,这里优先尝试 LibreOffice headless 转换;
            不可用时返回 ExpertError,提示用户安装 LibreOffice 或在 Windows 用 COM。
        """
        import shutil
        import subprocess

        err = self._validate_path(
            path, must_exist=True, allowed_exts=("pptx",)
        )
        if err:
            return self._failure(err)
        if not output_dir or not isinstance(output_dir, str):
            return self._failure("output_dir must be a non-empty string")

        try:
            # 优先 Windows COM(PowerPoint 已安装)
            if os.name == "nt":
                com_result = self._export_images_via_com(
                    path, output_dir, slide_indices
                )
                if com_result.success or com_result.metadata.get("tried_com"):
                    return com_result

            # LibreOffice headless
            soffice = shutil.which("soffice") or shutil.which("libreoffice")
            if soffice:
                abs_path = os.path.abspath(path)
                abs_out = os.path.abspath(output_dir)
                os.makedirs(abs_out, exist_ok=True)
                # 先转 PDF(LibreOffice 子进程,设超时避免挂死)
                try:
                    subprocess.run(
                        [soffice, "--headless", "--convert-to", "pdf",
                         "--outdir", abs_out, abs_path],
                        check=True, capture_output=True, timeout=120,
                    )
                except subprocess.CalledProcessError as e:
                    return self._failure(
                        f"LibreOffice convert-to pdf failed: "
                        f"{e.stderr.decode('utf-8', errors='replace') if e.stderr else str(e)}",
                    )
                except subprocess.TimeoutExpired:
                    return self._failure("LibreOffice convert-to pdf timeout (120s)")

                pdf_name = os.path.splitext(os.path.basename(abs_path))[0] + ".pdf"
                pdf_path = os.path.join(abs_out, pdf_name)
                if not os.path.exists(pdf_path):
                    return self._failure(f"LibreOffice produced no PDF: {pdf_path}")
                # 再用 PyMuPDF 把 PDF 转图片
                try:
                    import fitz  # PyMuPDF
                except ImportError:
                    return self._failure(
                        "export_images requires PyMuPDF (pip install pymupdf) "
                        "for image extraction",
                    )

                doc = None
                try:
                    doc = fitz.open(pdf_path)
                    total = doc.page_count
                    target_indices = slide_indices if slide_indices else list(range(1, total + 1))
                    image_paths: list[str] = []
                    for idx in target_indices:
                        if idx < 1 or idx > total:
                            continue
                        page = doc[idx - 1]
                        pix = page.get_pixmap(dpi=120)
                        out_file = os.path.join(abs_out, f"slide_{idx:03d}.png")
                        pix.save(out_file)
                        image_paths.append(out_file)
                    return self._success(
                        image_paths, count=len(image_paths), backend="libreoffice"
                    )
                except (PermissionError, IOError) as e:
                    return self._failure(f"export_images IO failed: {e}")
                except Exception as e:
                    return self._failure(f"export_images (libreoffice) failed: {e}")
                finally:
                    if doc is not None:
                        try:
                            doc.close()
                        except Exception:
                            pass

            return self._failure(
                "export_images requires either PowerPoint (Windows COM) or LibreOffice. "
                "Install one of them, or run on Windows with pywin32 installed.",
            )
        except Exception as e:
            return self._failure(f"export_images failed: {e}")

    # ------------------------------------------------------------------
    # 读取
    # ------------------------------------------------------------------

    def read(
        self,
        path: str,
        slide_range: Optional[tuple[int, int]] = None,
    ) -> ExpertResult:
        """读取 PPT 内容(每页标题/文本/图片/表格/备注)。

        Args:
            path: .pptx 路径
            slide_range: (start, end) 1-based 闭区间,限定读取的幻灯片范围;None 全部

        Returns:
            ExpertResult(output={
                slide_count, slides: [{index, title, text_blocks,
                has_image, has_table, tables, notes}], raw_text
            })
        """
        err = self._validate_path(
            path, must_exist=True, allowed_exts=("pptx",)
        )
        if err:
            return self._failure(err)
        if slide_range is not None:
            if (not isinstance(slide_range, (tuple, list))
                    or len(slide_range) != 2
                    or not all(isinstance(i, int) and not isinstance(i, bool)
                               for i in slide_range)):
                return self._failure(
                    "slide_range must be a tuple of two ints (start, end)"
                )
            if slide_range[0] < 1 or slide_range[1] < slide_range[0]:
                return self._failure(
                    "slide_range must satisfy 1 <= start <= end"
                )

        try:
            self._require_lib("pptx")
            from pptx import Presentation
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation(path)
            total = len(prs.slides)
            if slide_range is not None:
                start = max(1, slide_range[0])
                end = min(total, slide_range[1])
                if start > end:
                    return self._failure(
                        f"slide_range {slide_range} out of range (total={total})"
                    )
                target_indices = list(range(start - 1, end))
            else:
                target_indices = list(range(total))

            slides_data = []
            raw_parts: list[str] = []
            for i in target_indices:
                slide = prs.slides[i]
                # 标题 placeholder(slide.shapes.title 无则返回 None)
                # 注意:每次访问 shapes.title 返回新代理对象,is 比较失效,
                # 故用 shape_id 判等以跳过标题 shape
                title_shape = None
                title_shape_id = None
                try:
                    title_shape = slide.shapes.title
                    if title_shape is not None:
                        title_shape_id = title_shape.shape_id
                except Exception:
                    title_shape = None
                title = ""
                if title_shape is not None:
                    try:
                        title = (title_shape.text or "").strip()
                    except Exception:
                        title = ""

                text_blocks: list[str] = []
                has_image = False
                has_table = False
                tables: list[list[list[str]]] = []
                for shape in slide.shapes:
                    if title_shape_id is not None and getattr(shape, "shape_id", None) == title_shape_id:
                        continue  # 标题已单独提取
                    # 表格
                    if getattr(shape, "has_table", False):
                        has_table = True
                        try:
                            tbl = shape.table
                            rows = [
                                [cell.text for cell in row.cells]
                                for row in tbl.rows
                            ]
                            tables.append(rows)
                        except Exception:
                            pass
                        continue
                    # 图片(shape_type == 13 = PICTURE)
                    try:
                        if shape.shape_type == 13:
                            has_image = True
                            continue
                    except Exception:
                        pass
                    # 文本框
                    if shape.has_text_frame:
                        txt = shape.text_frame.text
                        if txt and txt.strip():
                            text_blocks.append(txt)

                # 备注(has_notes_slide 不创建新 notes slide)
                notes = ""
                try:
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text or ""
                except Exception:
                    notes = ""

                slides_data.append({
                    "index": i + 1,
                    "title": title,
                    "text_blocks": text_blocks,
                    "has_image": has_image,
                    "has_table": has_table,
                    "tables": tables,
                    "notes": notes,
                })
                if title:
                    raw_parts.append(title)
                raw_parts.extend(text_blocks)
                for tbl in tables:
                    for row in tbl:
                        raw_parts.append("\t".join(row))

            raw_text = "\n".join(raw_parts)
            return self._success(
                output={
                    "slide_count": total,
                    "slides": slides_data,
                    "raw_text": raw_text,
                },
                read_count=len(slides_data),
            )
        except (PermissionError, IOError) as e:
            return self._failure(f"read ppt IO failed: {e}")
        except Exception as e:
            return self._failure(f"read ppt failed: {e}")

    def extract_text(self, path: str) -> ExpertResult:
        """仅提取 PPT 纯文本(output=str)。

        Args:
            path: .pptx 路径

        Returns:
            ExpertResult(output=raw_text)
        """
        r = self.read(path)
        if not r.success:
            return r
        return self._success(
            r.output.get("raw_text", ""),
            slide_count=r.output.get("slide_count", 0),
        )

    def extract_slides(self, path: str) -> ExpertResult:
        """提取幻灯片结构(output=list[dict])。

        Args:
            path: .pptx 路径

        Returns:
            ExpertResult(output=[{index, title, text_blocks, has_image,
            has_table, tables, notes}, ...])
        """
        r = self.read(path)
        if not r.success:
            return r
        return self._success(
            r.output.get("slides", []),
            slide_count=r.output.get("slide_count", 0),
        )

    def _export_images_via_com(
        self,
        path: str,
        output_dir: str,
        slide_indices: Optional[list[int]],
    ) -> ExpertResult:
        """通过 Windows PowerPoint COM 接口导出图片(内部辅助)。

        确保 COM 调用后释放:Pres.Close + ppt_app.Quit + CoUninitialize。
        """
        try:
            import pythoncom
            import win32com.client  # type: ignore
        except ImportError:
            # 未装 pywin32,标记 tried_com=False 让上层走 LibreOffice
            return self._failure("pywin32 not installed", tried_com=False)

        pythoncom.CoInitialize()
        ppt_app = None
        pres = None
        try:
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            abs_path = os.path.abspath(path)
            abs_out = os.path.abspath(output_dir)
            os.makedirs(abs_out, exist_ok=True)
            pres = ppt_app.Presentations.Open(abs_path, WithWindow=False)
            total = pres.Slides.Count
            target_indices = slide_indices if slide_indices else list(range(1, total + 1))
            image_paths: list[str] = []
            for idx in target_indices:
                if idx < 1 or idx > total:
                    continue
                out_file = os.path.join(abs_out, f"slide_{idx:03d}.png")
                # PNG 导出,1280x720 适合预览
                pres.Slides(idx).Export(out_file, "PNG", 1280, 720)
                image_paths.append(out_file)
            return self._success(
                image_paths, count=len(image_paths), backend="com", tried_com=True
            )
        except Exception as e:
            return self._failure(
                f"export_images (COM) failed: {e}", tried_com=True
            )
        finally:
            # 严格按 COM 反向释放顺序:presentation → app → COM 库
            if pres is not None:
                try:
                    pres.Close()
                except Exception:
                    pass
            if ppt_app is not None:
                try:
                    ppt_app.Quit()
                except Exception:
                    pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass
