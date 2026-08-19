"""格式统一器(Phase 5.6)。

按规范统一 Word/Excel/PPT 文档格式,采用 Run 级原地修改策略,
保留主标题与代码块,对分类标题(第X部分/日期)加粗,正文不加粗。

参考实现:
  - Office-Word-MCP-Server/word_document_server/utils/document_utils.py
    (逐 run 操作,不整段重写)
  - Office-Word-MCP-Server/word_document_server/tools/format_tools.py
    (run 重建与格式应用)
  - fnixagent.office.word.WordExpert._set_run_font
    (rPr → rFonts → set w:eastAsia/w:ascii/w:hAnsi)

设计要点:
  - 设置东亚字体需同时操作 XML: rPr → rFonts → set w:eastAsia/w:ascii/w:hAnsi
  - 遍历 paragraph.runs 逐个设置,不整段重写,保留 run 边界
  - 跳过空 run(无文本),避免引入无意义 rPr
  - 返回 FormatReport 统计(total/modified/skipped + details)
"""

# -*- coding: utf-8 -*-
# Copyright (C) 2026 FnixAgent. All rights reserved.
# Software Name: FnixAgent 智能工作台系统 V1.0
# This software and its source code are proprietary and confidential.
# Unauthorized copying, modification, distribution, or use is strictly prohibited.

from __future__ import annotations

import os
import re
from dataclasses import asdict, dataclass, field
from typing import Any

from fnixagent.office.base import BaseExpert, ExpertError, ExpertResult

# 等宽代码字体集合(用于代码块识别)
_CODE_FONTS: frozenset[str] = frozenset(
    {
        "Consolas",
        "Courier New",
        "Courier",
        "Monaco",
        "Menlo",
        "Source Code Pro",
        "JetBrains Mono",
        "DejaVu Sans Mono",
        "Ubuntu Mono",
        "Fira Code",
        "Cascadia Code",
    }
)

# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------


@dataclass
class FormatSpec:
    """格式规范配置(可序列化)。

    Attributes:
        body_font: 正文中文字体名(如 "宋体")
        body_size_pt: 正文字号(磅,小四=12.0)
        body_bold: 正文是否加粗
        body_color: 正文颜色(十六进制如 "000000");None 不修改
        title_font: 分类标题中文字体名
        title_size_pt: 分类标题字号(磅)
        title_bold: 分类标题是否加粗
        main_title_keep: 主标题是否保持原样不统一
        heading_levels: 各级 Heading 样式配置 {level: (font, size_pt, bold)}
        normalize_whitespace: 是否合并多余空白(连续空白压成单个空格)
        preserve_code_blocks: 代码块是否保留等宽字体不统一
    """

    body_font: str = "宋体"
    body_size_pt: float = 12.0  # 小四
    body_bold: bool = False
    body_color: str | None = None  # 十六进制如 "000000"
    title_font: str = "宋体"
    title_size_pt: float = 12.0
    title_bold: bool = True
    main_title_keep: bool = True  # 主标题保持原样不统一
    heading_levels: dict[int, tuple] = field(
        default_factory=lambda: {
            1: ("黑体", 16, True),
            2: ("黑体", 14, True),
            3: ("黑体", 12, True),
        }
    )
    normalize_whitespace: bool = True  # 是否合并多余空白
    preserve_code_blocks: bool = True  # 代码块保留等宽字体


@dataclass
class FormatReport:
    """格式化结果报告。

    Attributes:
        total_runs: 处理的 run 总数(含跳过)
        modified_runs: 实际修改格式的 run 数
        skipped_runs: 跳过的 run 数(主标题/代码块/空 run)
        details: 每处修改/跳过的记录 dict 列表
    """

    total_runs: int = 0
    modified_runs: int = 0
    skipped_runs: int = 0
    details: list[dict] = field(default_factory=list)


# ---------------------------------------------------------------------------
# 格式统一器
# ---------------------------------------------------------------------------


class FormatNormalizer(BaseExpert):
    """格式统一器(Phase 5.6)。

    按 FormatSpec 统一 Word/Excel/PPT 文档格式,Run 级原地修改,
    保留主标题与代码块,分类标题(第X部分/日期)加粗,正文不加粗。

    能力边界:
      - 仅处理 .docx/.xlsx/.pptx(不支持旧版 .doc/.xls/.ppt)
      - 跨 run 文本不拼接(逐 run 替换,可能漏匹配跨 run 关键词)
      - 代码块识别基于等宽字体启发式,非语义级
    """

    @property
    def name(self) -> str:
        return "format_normalizer"

    # ------------------------------------------------------------------
    # 公共 API
    # ------------------------------------------------------------------

    def normalize(
        self,
        path: str,
        spec: FormatSpec | None = None,
        output_path: str | None = None,
    ) -> ExpertResult:
        """按扩展名派发到对应文档类型的统一方法。

        Args:
            path: 输入文档路径
            spec: 格式规范;None 使用默认 FormatSpec()
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": {...}})
        """
        ext = os.path.splitext(path)[1].lstrip(".").lower()
        if ext == "docx":
            return self.normalize_word(path, spec, output_path)
        if ext == "xlsx":
            return self.normalize_excel(path, spec, output_path)
        if ext == "pptx":
            return self.normalize_ppt(path, spec, output_path)
        return self._failure(f"unsupported extension: .{ext}, allowed: docx/xlsx/pptx")

    def normalize_word(
        self,
        path: str,
        spec: FormatSpec | None = None,
        output_path: str | None = None,
    ) -> ExpertResult:
        """统一 Word 文档格式。

        - 保留主标题(若 spec.main_title_keep)
        - 分类标题(第X部分/日期/小节标题)用 title_font/title_size/title_bold
        - Heading 1/2/3 样式段落用 heading_levels 配置
        - 正文用 body_font/body_size/body_bold,不加粗
        - 代码块(等宽字体)保留原字体

        Args:
            path: 输入 .docx 路径
            spec: 格式规范;None 使用默认
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": FormatReport dict})
        """
        spec = spec or FormatSpec()
        err = self._validate_path(path, must_exist=True, allowed_exts=("docx",))
        if err:
            return self._failure(err)
        if output_path:
            err = self._validate_path(output_path, allowed_exts=("docx",))
            if err:
                return self._failure(err)

        try:
            self._require_lib("docx")
            from docx import Document
        except ExpertError as e:
            return self._failure(str(e))

        try:
            doc = Document(path)
            report = FormatReport()

            # 处理正文段落
            self._normalize_paragraphs(doc.paragraphs, spec, report)

            # 处理表格内段落
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        self._normalize_paragraphs(cell.paragraphs, spec, report)

            save_path = output_path or path
            doc.save(save_path)
            return self._success(save_path, report=asdict(report))
        except (OSError, PermissionError) as e:
            return self._failure(f"normalize_word IO failed: {e}")
        except Exception as e:
            return self._failure(f"normalize_word failed: {e}")

    def normalize_excel(
        self,
        path: str,
        spec: FormatSpec | None = None,
        output_path: str | None = None,
    ) -> ExpertResult:
        """统一 Excel 字体字号。

        遍历所有 sheet 的非空单元格,应用 body_font/body_size/body_bold。
        保留单元格对齐/边框/填充/数字格式。

        Args:
            path: 输入 .xlsx 路径
            spec: 格式规范;None 使用默认
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": FormatReport dict})
        """
        spec = spec or FormatSpec()
        err = self._validate_path(path, must_exist=True, allowed_exts=("xlsx",))
        if err:
            return self._failure(err)
        if output_path:
            err = self._validate_path(output_path, allowed_exts=("xlsx",))
            if err:
                return self._failure(err)

        try:
            self._require_lib("openpyxl")
            from openpyxl import load_workbook
            from openpyxl.styles import Font
        except ExpertError as e:
            return self._failure(str(e))

        try:
            wb = load_workbook(path)
            report = FormatReport()

            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        report.total_runs += 1
                        # 保留原对齐/边框/填充/数字格式
                        old_align = cell.alignment
                        old_border = cell.border
                        old_fill = cell.fill
                        old_num_fmt = cell.number_format

                        cell.font = Font(
                            name=spec.body_font,
                            size=spec.body_size_pt,
                            bold=spec.body_bold,
                        )
                        # 恢复非字体格式
                        cell.alignment = old_align
                        cell.border = old_border
                        cell.fill = old_fill
                        cell.number_format = old_num_fmt
                        report.modified_runs += 1

            save_path = output_path or path
            wb.save(save_path)
            return self._success(save_path, report=asdict(report))
        except (OSError, PermissionError) as e:
            return self._failure(f"normalize_excel IO failed: {e}")
        except Exception as e:
            return self._failure(f"normalize_excel failed: {e}")

    def normalize_ppt(
        self,
        path: str,
        spec: FormatSpec | None = None,
        output_path: str | None = None,
    ) -> ExpertResult:
        """统一 PPT 字体。

        遍历所有 slide 的文本框,逐 run 应用 body_font/body_size/body_bold。
        跳过空 run。

        Args:
            path: 输入 .pptx 路径
            spec: 格式规范;None 使用默认
            output_path: 输出路径;None 覆盖原文件

        Returns:
            ExpertResult(output=save_path, metadata={"report": FormatReport dict})
        """
        spec = spec or FormatSpec()
        err = self._validate_path(path, must_exist=True, allowed_exts=("pptx",))
        if err:
            return self._failure(err)
        if output_path:
            err = self._validate_path(output_path, allowed_exts=("pptx",))
            if err:
                return self._failure(err)

        try:
            self._require_lib("pptx")
            from pptx import Presentation
            from pptx.util import Pt
        except ExpertError as e:
            return self._failure(str(e))

        try:
            prs = Presentation(path)
            report = FormatReport()

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            report.total_runs += 1
                            if not run.text:
                                report.skipped_runs += 1
                                continue
                            # 合并空白
                            if spec.normalize_whitespace:
                                new_text = re.sub(r"\s+", " ", run.text)
                                if new_text != run.text:
                                    run.text = new_text
                            # 应用字体
                            run.font.name = spec.body_font
                            run.font.size = Pt(spec.body_size_pt)
                            run.font.bold = spec.body_bold
                            if spec.body_color:
                                self._apply_ppt_color(run, spec.body_color)
                            report.modified_runs += 1

            save_path = output_path or path
            prs.save(save_path)
            return self._success(save_path, report=asdict(report))
        except (OSError, PermissionError) as e:
            return self._failure(f"normalize_ppt IO failed: {e}")
        except Exception as e:
            return self._failure(f"normalize_ppt failed: {e}")

    # ------------------------------------------------------------------
    # Run 级格式应用
    # ------------------------------------------------------------------

    def apply_to_run(
        self,
        run: Any,
        font: str,
        size_pt: float,
        bold: bool,
        color: str | None = None,
    ) -> None:
        """应用格式到单个 run(设置 font.name + eastAsia + size + bold + color)。

        参考实现: fnixagent.office.word.WordExpert._set_run_font
        设置东亚字体需同时操作 XML:
          rPr → rFonts → set w:eastAsia/w:ascii/w:hAnsi/w:cs

        Args:
            run: python-docx 的 Run 对象
            font: 中文字体名(如 "宋体")
            size_pt: 字号(磅)
            bold: 是否加粗
            color: 十六进制颜色(如 "000000" 或 "#000000");None 不修改
        """
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls, qn
        from docx.shared import Pt, RGBColor

        # ASCII 字体设为 Times New Roman,中文用 font 参数
        run.font.name = "Times New Roman"
        run.font.size = Pt(size_pt)
        run.font.bold = bold

        # 颜色
        if color:
            c = color.lstrip("#").upper()
            if len(c) == 6:
                try:
                    run.font.color.rgb = RGBColor(
                        int(c[0:2], 16),
                        int(c[2:4], 16),
                        int(c[4:6], 16),
                    )
                except ValueError:
                    pass

        # 东亚字体:操作 rPr → rFonts
        rpr = run._element.get_or_add_rPr()
        rfonts = rpr.find(qn("w:rFonts"))
        if rfonts is None:
            rfonts = parse_xml(
                f"<w:rFonts {nsdecls('w')} "
                f'w:ascii="Times New Roman" w:hAnsi="Times New Roman" '
                f'w:eastAsia="{font}" w:cs="Times New Roman"/>'
            )
            rpr.insert(0, rfonts)
        else:
            rfonts.set(qn("w:ascii"), "Times New Roman")
            rfonts.set(qn("w:hAnsi"), "Times New Roman")
            rfonts.set(qn("w:eastAsia"), font)
            rfonts.set(qn("w:cs"), "Times New Roman")

    # ------------------------------------------------------------------
    # 标题识别
    # ------------------------------------------------------------------

    def is_main_title(self, text: str) -> bool:
        """判断是否主标题(保持原样)。

        启发式规则:
          - 文本非空且较短(<= 50 字符)
          - 无句末标点(. 。 ! ? ！ ？ ; ；)
          - 不含分类标题特征(第X部分/中文序号/阿拉伯序号)
          - 短文本且不含逗号(避免普通句子)

        Args:
            text: 段落文本

        Returns:
            True 视为主标题
        """
        if not text:
            return False
        t = text.strip()
        if not t or len(t) > 50:
            return False
        # 句末标点 → 不是主标题
        if t.endswith(("。", ".", "!", "?", "！", "？", ";", "；")):
            return False
        # 分类标题特征 → 不是主标题(交给 is_title_paragraph)
        if re.match(r"^第[一二三四五六七八九十百零\d]+[部分章节条条款]", t):
            return False
        if re.match(r"^[一二三四五六七八九十]+、", t):
            return False
        if re.match(r"^\d+[\.．、]\s*\S", t):
            return False
        # 日期格式 → 不是主标题(交给 is_title_paragraph)
        if re.match(r"^\d{4}[-年]\d{1,2}[-月]\d{1,2}日?$", t):
            return False
        # 选项段落(A./B./C./D. 开头)→ 不是主标题
        if re.match(r"^[A-Z]\.\s", t):
            return False
        # 题干(【单选题】等题型标签)→ 不是主标题
        if re.match(r"^【.+】", t):
            return False
        # 答案行(【答案】)→ 不是主标题
        if "【答案】" in t:
            return False
        # 短文本且无逗号 → 视为主标题
        if len(t) <= 30 and ("，" not in t and "," not in t):
            return True
        return False

    def is_title_paragraph(self, text: str) -> bool:
        """判断是否标题性段落(第X部分/日期/小节标题)。

        匹配模式:
          - 第X部分/章/节/条/款
          - 中文序号开头(一、二、三、)
          - 阿拉伯数字序号(1. 2. 3.)
          - 日期格式(2025年X月X日 / 2025-01-01)
          - 短文本(< 30 字)无句末标点且无逗号

        Args:
            text: 段落文本

        Returns:
            True 视为标题性段落
        """
        if not text:
            return False
        t = text.strip()
        if not t:
            return False
        # 第X部分/章/节
        if re.match(r"^第[一二三四五六七八九十百零\d]+[部分章节条条款]", t):
            return True
        # 中文序号开头
        if re.match(r"^[一二三四五六七八九十]+、", t):
            return True
        # 阿拉伯数字序号
        if re.match(r"^\d+[\.．、]\s*\S", t):
            return True
        # 日期格式
        if re.match(r"^\d{4}[-年]\d{1,2}[-月]\d{1,2}日?$", t):
            return True
        # 排除选项段落(A./B./C./D. 开头)、题干、答案行
        if re.match(r"^[A-Z]\.\s", t):
            return False
        if re.match(r"^【.+】", t):
            return False
        if "【答案】" in t:
            return False
        # 短标题启发式
        if (
            len(t) < 30
            and not t.endswith(("。", ".", "!", "?", "！", "？", ";", "；"))
            and ("，" not in t and "," not in t)
        ):
            return True
        return False

    # ------------------------------------------------------------------
    # 内部工具
    # ------------------------------------------------------------------

    def _normalize_paragraphs(
        self,
        paragraphs: Any,
        spec: FormatSpec,
        report: FormatReport,
    ) -> None:
        """对一组段落应用格式统一,结果累计到 report。"""
        for para in paragraphs:
            text = para.text.strip() if para.text else ""
            if not text:
                continue

            style_name = ""
            if para.style and para.style.name:
                style_name = para.style.name

            # 主标题跳过
            if spec.main_title_keep:
                if style_name == "Title" or self.is_main_title(text):
                    run_count = len(para.runs)
                    report.total_runs += run_count
                    report.skipped_runs += run_count
                    report.details.append(
                        {
                            "type": "skip_main_title",
                            "text": text[:50],
                            "runs": run_count,
                        }
                    )
                    continue

            # 决定字体/字号/加粗
            font = spec.body_font
            size = spec.body_size_pt
            bold = spec.body_bold
            color = spec.body_color
            is_title = False

            if style_name == "Heading 1" and 1 in spec.heading_levels:
                h = spec.heading_levels[1]
                font, size, bold = h[0], h[1], h[2]
                is_title = True
            elif style_name == "Heading 2" and 2 in spec.heading_levels:
                h = spec.heading_levels[2]
                font, size, bold = h[0], h[1], h[2]
                is_title = True
            elif style_name == "Heading 3" and 3 in spec.heading_levels:
                h = spec.heading_levels[3]
                font, size, bold = h[0], h[1], h[2]
                is_title = True
            elif self.is_title_paragraph(text):
                font = spec.title_font
                size = spec.title_size_pt
                bold = spec.title_bold
                is_title = True

            # 逐 run 处理
            for run in para.runs:
                report.total_runs += 1
                if not run.text:
                    report.skipped_runs += 1
                    continue
                # 代码块保留
                if spec.preserve_code_blocks and self._is_code_run(run):
                    report.skipped_runs += 1
                    report.details.append(
                        {
                            "type": "skip_code",
                            "text": run.text[:50],
                        }
                    )
                    continue
                # 合并空白
                if spec.normalize_whitespace:
                    new_text = re.sub(r"\s+", " ", run.text)
                    if new_text != run.text:
                        run.text = new_text
                # 应用格式
                self.apply_to_run(run, font, size, bold, color)
                report.modified_runs += 1
                report.details.append(
                    {
                        "type": "modify",
                        "is_title": is_title,
                        "font": font,
                        "size": size,
                        "bold": bold,
                    }
                )

    @staticmethod
    def _is_code_run(run: Any) -> bool:
        """检测 run 是否为代码字体(等宽)。

        Args:
            run: python-docx Run 对象

        Returns:
            True 视为代码 run
        """
        try:
            from docx.oxml.ns import qn

            # 检查 font.name 属性
            font_name = run.font.name
            if font_name and font_name in _CODE_FONTS:
                return True
            # 检查 XML 中的 ascii/hAnsi 字体
            rpr = run._element.find(qn("w:rPr"))
            if rpr is not None:
                rfonts = rpr.find(qn("w:rFonts"))
                if rfonts is not None:
                    for attr in ("w:ascii", "w:hAnsi", "w:eastAsia"):
                        val = rfonts.get(qn(attr))
                        if val and val in _CODE_FONTS:
                            return True
        except Exception:
            pass
        return False

    @staticmethod
    def _apply_ppt_color(run: Any, color: str) -> None:
        """为 PPT run 设置字体颜色。

        Args:
            run: python-pptx Run 对象
            color: 十六进制颜色(如 "FF0000" 或 "#FF0000")
        """
        try:
            from pptx.dml.color import RGBColor

            c = color.lstrip("#").upper()
            if len(c) == 6:
                run.font.color.rgb = RGBColor(
                    int(c[0:2], 16),
                    int(c[2:4], 16),
                    int(c[4:6], 16),
                )
        except Exception:
            pass
