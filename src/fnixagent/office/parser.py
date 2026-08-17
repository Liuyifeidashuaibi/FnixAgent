"""Parser Expert(P2-9)。

文档解析:把 Word/Excel/PDF/HTML 等解析为结构化数据(段落/表格/表单)。

专家职责:
  - 按扩展名自动派发解析器(docx/xlsx/pdf/html/txt/csv/json)
  - 表格抽取、表单字段抽取(正则)、布局检测
  - 统一 Element 模型输出

底层依赖:
  - python-docx / openpyxl / pypdf / pdfplumber / beautifulsoup4(按需)

降级策略:
  - 依赖缺失 → ExpertError 提示安装
  - 工作簿 read_only 模式,close 在 finally 中确保释放
  - 路径/扩展名校验前置,避免无效 IO

设计思路(Unstructured.io):
  - 统一 Element 模型: 所有解析产物都是 Element 子类实例,具备
    category/text/metadata/to_dict() 接口,方便下游统一处理。
  - 自描述 FileType 枚举: 每个枚举值携带支持的扩展名、能力标志
    (supports_tables/supports_metadata),由 from_extension() 派发。
  - 装饰器栈后处理: @apply_metadata / @add_chunking_strategy 在
    parser 返回 Element 列表后统一应用元数据与分块策略,parser
    本身只关心"抽取元素"。
  - ParseOptions: 替代散落的布尔参数,集中表达解析选项。
  - cached_property: word_count/char_count 等派生属性惰性计算。
"""

# -*- coding: utf-8 -*-
# Copyright (C) 2026 FnixAgent. All rights reserved.
# Software Name: FnixAgent 智能工作台系统 V1.0
# This software and its source code are proprietary and confidential.
# Unauthorized copying, modification, distribution, or use is strictly prohibited.

from __future__ import annotations

import functools
import os
import re
from collections.abc import Callable
from dataclasses import dataclass, field
from enum import Enum
from functools import cached_property
from typing import Any

from fnixagent.office.base import BaseExpert, ExpertError, ExpertResult

# ---------------------------------------------------------------------------
# 统一 Element 模型
# ---------------------------------------------------------------------------

@dataclass
class Element:
    """所有解析产物的基础元素。

    Attributes:
        text: 元素文本内容(表格为序列化文本)
        metadata: 元素元数据(页码/样式/位置/来源 等)
    """

    text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def category(self) -> str:
        """元素类别(子类覆盖)。"""
        return "UncategorizedText"

    @cached_property
    def word_count(self) -> int:
        """词数(按空白切分)。"""
        return len(self.text.split())

    @cached_property
    def char_count(self) -> int:
        """字符数。"""
        return len(self.text)

    def to_dict(self) -> dict[str, Any]:
        """序列化为字典(供 JSON 输出/Agent 工具结果)。"""
        return {
            "type": self.category,
            "text": self.text,
            "metadata": dict(self.metadata),
        }

    def __eq__(self, other: object) -> bool:
        if not isinstance(other, Element):
            return NotImplemented
        return self.category == other.category and self.text == other.text

    def __hash__(self) -> int:
        return hash((self.category, self.text))

class Title(Element):
    """标题元素(Heading 1/2/3 等)。"""

    @property
    def category(self) -> str:
        return "Title"

class NarrativeText(Element):
    """正文叙述文本。"""

    @property
    def category(self) -> str:
        return "NarrativeText"

class ListItem(Element):
    """列表项。"""

    @property
    def category(self) -> str:
        return "ListItem"

class Table(Element):
    """表格元素。

    Attributes:
        rows: 二维行数据(首行为表头)
    """

    rows: list[list[str]] = field(default_factory=list)

    def __init__(
        self, rows: list[list[str]] | None = None, text: str = "", metadata: dict | None = None
    ) -> None:
        super().__init__(text=text or "", metadata=metadata or {})
        self.rows = rows or []
        # 自动序列化 text(便于全文检索)
        if not self.text and self.rows:
            self.text = "\n".join("\t".join(str(c) for c in r) for r in self.rows)

    @property
    def category(self) -> str:
        return "Table"

    @cached_property
    def row_count(self) -> int:
        return len(self.rows)

    @cached_property
    def column_count(self) -> int:
        return max((len(r) for r in self.rows), default=0)

    def to_dict(self) -> dict[str, Any]:
        d = super().to_dict()
        d["rows"] = self.rows
        d["row_count"] = self.row_count
        d["column_count"] = self.column_count
        return d

class Image(Element):
    """图片元素(占位,不抽取图像内容)。"""

    @property
    def category(self) -> str:
        return "Image"

class Header(Element):
    """页眉元素。"""

    @property
    def category(self) -> str:
        return "Header"

class Footer(Element):
    """页脚元素。"""

    @property
    def category(self) -> str:
        return "Footer"

class PageBreak(Element):
    """分页符元素(text 为空)。"""

    def __init__(self, metadata: dict | None = None) -> None:
        super().__init__(text="", metadata=metadata or {})

    @property
    def category(self) -> str:
        return "PageBreak"

# ---------------------------------------------------------------------------
# 自描述 FileType 枚举
# ---------------------------------------------------------------------------

class FileType(Enum):
    """文件类型枚举,每个值携带能力描述与扩展名映射。

    通过 from_extension() 由扩展名派发到对应枚举值,再由枚举值的能力
    标志决定是否支持表格抽取/元数据抽取,避免散落的 if-elif 链。
    """

    DOCX = "docx"  # Word
    XLSX = "xlsx"  # Excel
    PDF = "pdf"
    HTML = "html"
    TEXT = "text"  # txt/md
    CSV = "csv"
    JSON = "json"
    PPTX = "pptx"  # PowerPoint
    UNKNOWN = "unknown"

    # -- 能力标志(自描述) -------------------------------------------------
    @property
    def supports_tables(self) -> bool:
        """该类型是否原生支持表格抽取。"""
        return self in (
            FileType.DOCX,
            FileType.XLSX,
            FileType.HTML,
            FileType.CSV,
            FileType.PDF,
            FileType.PPTX,
        )

    @property
    def supports_metadata(self) -> bool:
        """该类型是否支持元数据抽取。"""
        return self in (FileType.DOCX, FileType.XLSX, FileType.PDF, FileType.HTML, FileType.PPTX)

    @property
    def extensions(self) -> tuple[str, ...]:
        """该类型对应的扩展名(小写无点)。"""
        return {
            FileType.DOCX: ("docx", "doc"),
            FileType.XLSX: ("xlsx", "xls"),
            FileType.PDF: ("pdf",),
            FileType.HTML: ("html", "htm"),
            FileType.TEXT: ("txt", "md"),
            FileType.CSV: ("csv",),
            FileType.JSON: ("json",),
            FileType.PPTX: ("pptx",),
            FileType.UNKNOWN: (),
        }[self]

    @classmethod
    def from_extension(cls, ext: str) -> FileType:
        """由扩展名派发到枚举值;未识别返回 UNKNOWN。

        Args:
            ext: 扩展名(可带点或不带点,大小写不敏感)

        Returns:
            对应的 FileType 枚举值
        """
        norm = ext.lstrip(".").lower()
        for ft in cls:
            if norm in ft.extensions:
                return ft
        return cls.UNKNOWN

# ---------------------------------------------------------------------------
# ParseOptions(替代散落的布尔参数)
# ---------------------------------------------------------------------------

@dataclass
class ParseOptions:
    """解析选项集中表达。

    Attributes:
        extract_tables: 是否抽取表格
        extract_metadata: 是否抽取元数据
        include_images: 是否包含图片元素(占位)
        chunking_strategy: 分块策略(None/"basic"/"by_page")
        chunk_size: 分块大小(字符数,仅 chunking_strategy 非空时生效)
    """

    extract_tables: bool = True
    extract_metadata: bool = True
    include_images: bool = False
    chunking_strategy: str | None = None
    chunk_size: int = 1200

# ---------------------------------------------------------------------------
# 装饰器栈
# ---------------------------------------------------------------------------

def apply_metadata(func: Callable) -> Callable:
    """装饰器:为 parser 返回的 Element 列表统一注入文件级元数据。

    被装饰的解析方法签名: (self, path, options) -> list[Element]
    装饰后: 自动为每个 Element.metadata 注入 file_path/file_type/source。
    """

    @functools.wraps(func)
    def wrapper(self: ParserExpert, path: str, options: ParseOptions) -> list[Element]:
        elements = func(self, path, options)
        file_type = os.path.splitext(path)[1].lstrip(".").lower()
        for el in elements:
            el.metadata.setdefault("file_path", path)
            el.metadata.setdefault("file_type", file_type)
            el.metadata.setdefault("source", os.path.basename(path))
        return elements

    return wrapper

def add_chunking_strategy(func: Callable) -> Callable:
    """装饰器:按 options.chunking_strategy 对 Element 列表分块。

    策略:
      - None: 不分块,原样返回
      - "basic": 按 chunk_size 字符数合并相邻 NarrativeText,保持 Title 完整
      - "by_page": 按 metadata.page_number 分组,跨页合并同类型元素
    """

    @functools.wraps(func)
    def wrapper(self: ParserExpert, path: str, options: ParseOptions) -> list[Element]:
        elements = func(self, path, options)
        if options.chunking_strategy is None:
            return elements
        if options.chunking_strategy == "basic":
            return _chunk_basic(elements, options.chunk_size)
        elif options.chunking_strategy == "by_page":
            return _chunk_by_page(elements)
        return elements

    return wrapper

def _chunk_basic(elements: list[Element], chunk_size: int) -> list[Element]:
    """basic 分块:合并相邻 NarrativeText 直至达到 chunk_size。"""
    if chunk_size <= 0:
        return elements
    result: list[Element] = []
    buffer: list[str] = []
    buffer_size = 0
    for el in elements:
        if isinstance(el, NarrativeText):
            if buffer_size + len(el.text) > chunk_size and buffer:
                result.append(
                    NarrativeText(
                        text="\n\n".join(buffer),
                        metadata={"chunk_index": len(result)},
                    )
                )
                buffer = []
                buffer_size = 0
            buffer.append(el.text)
            buffer_size += len(el.text)
        else:
            if buffer:
                result.append(
                    NarrativeText(
                        text="\n\n".join(buffer),
                        metadata={"chunk_index": len(result)},
                    )
                )
                buffer = []
                buffer_size = 0
            result.append(el)
    if buffer:
        result.append(
            NarrativeText(
                text="\n\n".join(buffer),
                metadata={"chunk_index": len(result)},
            )
        )
    return result

def _chunk_by_page(elements: list[Element]) -> list[Element]:
    """by_page 分块:按 page_number 分组合并同类型元素。"""
    grouped: dict[Any, list[Element]] = {}
    order: list[Any] = []
    for el in elements:
        page = el.metadata.get("page_number", 0)
        if page not in grouped:
            grouped[page] = []
            order.append(page)
        grouped[page].append(el)
    result: list[Element] = []
    for page in order:
        result.extend(grouped[page])
    return result

# ---------------------------------------------------------------------------
# ParserExpert
# ---------------------------------------------------------------------------

class ParserExpert(BaseExpert):
    """文档解析专家。

    全部方法返回 ExpertResult,output 为结构化数据。

    提供两套 API:
      - parse()(向后兼容):返回 dict 结构(type/paragraphs/tables/raw_text)
      - parse_elements()(P2-9 新增,):返回 Element 列表,
        支持 FileType 自描述派发、装饰器栈后处理、ParseOptions 集中配置

    能力边界:
      - PDF 表格抽取依赖 pdfplumber
      - 表单字段抽取为正则匹配,非语义理解
      - 布局检测为启发式(标题样式/表格/图片存在性)
    """

    @property
    def name(self) -> str:
        return "parser"

    # ------------------------------------------------------------------
    # 统一入口(向后兼容: dict 输出)
    # ------------------------------------------------------------------

    def parse(
        self,
        path: str,
        extract_tables: bool = True,
        extract_metadata: bool = True,
    ) -> ExpertResult:
        """根据扩展名自动选择解析器,返回结构化数据(dict 输出,向后兼容)。

        Args:
            path: 文件路径
            extract_tables: 是否抽取表格
            extract_metadata: 是否抽取 metadata

        Returns:
            ExpertResult(output={
                type, metadata, paragraphs, tables, raw_text
            })
        """
        err = self._validate_path(path, must_exist=True)
        if err:
            return self._failure(err)
        # P2-9: 用自描述 FileType 枚举派发(替代硬编码 dispatch dict)
        ft = FileType.from_extension(os.path.splitext(path)[1])
        if ft is FileType.UNKNOWN:
            ext = os.path.splitext(path)[1].lstrip(".").lower()
            return self._failure(f"unsupported file type: .{ext}")
        dispatch = {
            FileType.DOCX: self._parse_docx,
            FileType.XLSX: self._parse_xlsx,
            FileType.PDF: self._parse_pdf,
            FileType.HTML: self._parse_html,
            FileType.TEXT: self._parse_text,
            FileType.CSV: self._parse_csv,
            FileType.JSON: self._parse_json,
            FileType.PPTX: self._parse_pptx,
        }
        handler = dispatch.get(ft)
        if not handler:
            return self._failure(f"unsupported file type: {ft.value}")
        try:
            result = handler(path, extract_tables, extract_metadata)
            result.metadata["file_type"] = ft.value
            return result
        except (OSError, PermissionError) as e:
            return self._failure(f"parse IO failed: {e}")
        except Exception as e:
            return self._failure(f"parse failed: {e}")

    # ------------------------------------------------------------------
    # 统一入口(P2-9 新增: Element 列表输出,)
    # ------------------------------------------------------------------

    def parse_elements(
        self,
        path: str,
        options: ParseOptions | None = None,
    ) -> ExpertResult:
        """根据扩展名自动选择解析器,返回 Element 列表(Unstructured 风格)。

        与 parse() 的区别:
          - 输出为 Element 列表(Title/NarrativeText/Table/...),而非 dict
          - 支持 ParseOptions 集中配置(extract_tables/metadata/chunking 等)
          - 装饰器栈自动注入 file_path/file_type/source 元数据
          - 支持 chunking_strategy(basic/by_page)分块

        Args:
            path: 文件路径
            options: 解析选项;None 使用默认(全抽取,不分块)

        Returns:
            ExpertResult(output={
                "elements": [Element.to_dict(), ...],
                "file_type": str,
                "element_count": int,
                "metadata": dict,
            })
        """
        err = self._validate_path(path, must_exist=True)
        if err:
            return self._failure(err)
        opts = options or ParseOptions()
        ft = FileType.from_extension(os.path.splitext(path)[1])
        if ft is FileType.UNKNOWN:
            ext = os.path.splitext(path)[1].lstrip(".").lower()
            return self._failure(f"unsupported file type: .{ext}")
        dispatch = {
            FileType.DOCX: self._parse_docx_elements,
            FileType.XLSX: self._parse_xlsx_elements,
            FileType.PDF: self._parse_pdf_elements,
            FileType.HTML: self._parse_html_elements,
            FileType.TEXT: self._parse_text_elements,
            FileType.CSV: self._parse_csv_elements,
            FileType.JSON: self._parse_json_elements,
            FileType.PPTX: self._parse_pptx_elements,
        }
        handler = dispatch.get(ft)
        if not handler:
            return self._failure(f"unsupported file type: {ft.value}")
        try:
            elements = handler(path, opts)
            # 兼容能力标志:不支持表格时过滤掉 Table 元素
            if not opts.extract_tables:
                elements = [e for e in elements if not isinstance(e, Table)]
            if not opts.include_images:
                elements = [e for e in elements if not isinstance(e, Image)]
            metadata = {"file_type": ft.value, "source": os.path.basename(path)}
            return self._success(
                output={
                    "elements": [e.to_dict() for e in elements],
                    "file_type": ft.value,
                    "element_count": len(elements),
                    "metadata": metadata,
                },
                element_count=len(elements),
            )
        except (OSError, PermissionError) as e:
            return self._failure(f"parse_elements IO failed: {e}")
        except Exception as e:
            return self._failure(f"parse_elements failed: {e}")

    @staticmethod
    def elements_to_text(elements: list[Element], sep: str = "\n\n") -> str:
        """将 Element 列表拼接为纯文本(便于喂给 LLM)。

        Args:
            elements: Element 实例列表
            sep: 元素间分隔符

        Returns:
            拼接后的纯文本(PageBreak 元素跳过)
        """
        return sep.join(e.text for e in elements if not isinstance(e, PageBreak))

    @staticmethod
    def filter_by_category(elements: list[Element], category: str) -> list[Element]:
        """按类别过滤元素(如 "Title"/"Table"/"NarrativeText")。

        Args:
            elements: Element 实例列表
            category: 元素类别(Element.category 属性值)

        Returns:
            匹配类别的子列表
        """
        return [e for e in elements if e.category == category]

    # ------------------------------------------------------------------
    # 表格抽取
    # ------------------------------------------------------------------

    def parse_table(
        self,
        path: str,
        sheet_name: str | None = None,
        table_index: int = 0,
    ) -> ExpertResult:
        """专门抽取表格数据。

        Args:
            path: 文件路径
            sheet_name: Excel sheet 名(仅 xlsx 有效)
            table_index: 第几张表(0-based,Word/PDF 多表格场景)

        Returns:
            ExpertResult(output={headers, rows})
        """
        err = self._validate_path(path, must_exist=True)
        if err:
            return self._failure(err)
        err = self._validate_int(table_index, "table_index", min_value=0)
        if err:
            return self._failure(err)
        ext = os.path.splitext(path)[1].lstrip(".").lower()
        try:
            if ext in ("xlsx", "xls"):
                return self._parse_xlsx_table(path, sheet_name)
            elif ext in ("docx", "doc"):
                return self._parse_docx_table(path, table_index)
            elif ext == "csv":
                return self._parse_csv_table(path)
            elif ext == "pdf":
                return self._parse_pdf_table(path, table_index)
            else:
                return self._failure(f"table extraction not supported for .{ext}")
        except (OSError, PermissionError) as e:
            return self._failure(f"parse_table IO failed: {e}")
        except Exception as e:
            return self._failure(f"parse_table failed: {e}")

    # ------------------------------------------------------------------
    # 表单抽取
    # ------------------------------------------------------------------

    def parse_form(
        self,
        path: str,
        field_patterns: dict[str, str] | None = None,
    ) -> ExpertResult:
        """抽取表单字段(基于正则匹配键值对)。

        Args:
            path: 文件路径(支持 docx/pdf/txt)
            field_patterns: {field_name: regex};None 用内置常见字段

        Returns:
            ExpertResult(output={field: value})
        """
        err = self._validate_path(path, must_exist=True)
        if err:
            return self._failure(err)
        # 默认常见表单字段
        if field_patterns is None:
            field_patterns = {
                "name": r"(?:姓名|Name|名称)[:：\s]*([^\s,，。.]+)",
                "id_card": r"(?:身份证|ID)[:：\s]*([0-9Xx]{15,18})",
                "phone": r"(?:电话|手机|Phone|Tel)[:：\s]*([0-9\-\+\s]{7,20})",
                "email": r"(?:邮箱|Email|E-mail)[:：\s]*([\w\.\-]+@[\w\.\-]+)",
                "address": r"(?:地址|Address)[:：\s]*([^\n,，。.]+)",
                "date": r"(?:日期|Date)[:：\s]*([\d\-/年月日]+)",
                "amount": r"(?:金额|Amount)[:：\s]*([\d,，.]+)",
            }

        # 先抽取全文文本
        text_result = self._extract_plain_text(path)
        if not text_result.success:
            return text_result
        text = text_result.output

        fields: dict[str, str | None] = {}
        for fname, pattern in field_patterns.items():
            try:
                match = re.search(pattern, text)
                fields[fname] = match.group(1).strip() if match else None
            except re.error as e:
                return self._failure(f"invalid regex for field '{fname}': {e}")
        return self._success(fields, fields_found=sum(1 for v in fields.values() if v))

    # ------------------------------------------------------------------
    # 布局检测
    # ------------------------------------------------------------------

    def detect_layout(
        self,
        path: str,
    ) -> ExpertResult:
        """检测文档布局类型(段落/表格/标题/列表/图片的分布)。

        Args:
            path: 文件路径

        Returns:
            ExpertResult(output={
                layout_type, has_table, has_image, has_header,
                sections, paragraph_count, table_count
            })
        """
        err = self._validate_path(path, must_exist=True)
        if err:
            return self._failure(err)
        ext = os.path.splitext(path)[1].lstrip(".").lower()
        try:
            if ext in ("docx", "doc"):
                return self._detect_docx_layout(path)
            elif ext in ("xlsx", "xls"):
                return self._detect_xlsx_layout(path)
            elif ext == "pdf":
                return self._detect_pdf_layout(path)
            elif ext in ("html", "htm"):
                return self._detect_html_layout(path)
            elif ext == "pptx":
                return self._detect_pptx_layout(path)
            else:
                # 纯文本类型
                return self._success(
                    {
                        "layout_type": "plain_text",
                        "has_table": False,
                        "has_image": False,
                        "has_header": False,
                        "sections": 1,
                        "paragraph_count": 1,
                        "table_count": 0,
                    }
                )
        except (OSError, PermissionError) as e:
            return self._failure(f"detect_layout IO failed: {e}")
        except Exception as e:
            return self._failure(f"detect_layout failed: {e}")

    # ------------------------------------------------------------------
    # 内部解析器:docx
    # ------------------------------------------------------------------

    def _parse_docx(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        try:
            self._require_lib("docx")
            from docx import Document
        except ExpertError as e:
            return self._failure(str(e))

        doc = Document(path)
        paragraphs = []
        for p in doc.paragraphs:
            if p.text.strip():
                paragraphs.append(
                    {
                        "text": p.text,
                        "style": p.style.name if p.style else "Normal",
                    }
                )
        tables = []
        if extract_tables:
            for tbl in doc.tables:
                rows = []
                for row in tbl.rows:
                    rows.append([cell.text for cell in row.cells])
                tables.append(rows)
        metadata = {}
        if extract_metadata:
            cp = doc.core_properties
            metadata = {
                "title": cp.title or "",
                "author": cp.author or "",
                "subject": cp.subject or "",
                "created": str(cp.created) if cp.created else "",
                "modified": str(cp.modified) if cp.modified else "",
            }
        raw_text = "\n".join(p["text"] for p in paragraphs)
        return self._success(
            output={
                "type": "docx",
                "metadata": metadata,
                "paragraphs": paragraphs,
                "tables": tables,
                "raw_text": raw_text,
            },
            paragraph_count=len(paragraphs),
            table_count=len(tables),
        )

    def _parse_docx_table(self, path: str, table_index: int) -> ExpertResult:
        try:
            self._require_lib("docx")
            from docx import Document
        except ExpertError as e:
            return self._failure(str(e))
        doc = Document(path)
        if table_index < 0 or table_index >= len(doc.tables):
            return self._failure(f"table_index out of range: {table_index}")
        tbl = doc.tables[table_index]
        rows = [[cell.text for cell in row.cells] for row in tbl.rows]
        headers = rows[0] if rows else []
        return self._success(
            output={"headers": headers, "rows": rows[1:] if len(rows) > 1 else []},
            rows=len(rows),
        )

    def _detect_docx_layout(self, path: str) -> ExpertResult:
        try:
            self._require_lib("docx")
            from docx import Document
        except ExpertError as e:
            return self._failure(str(e))
        doc = Document(path)
        has_header = any(p.style and "Heading" in p.style.name for p in doc.paragraphs)
        # 图片检测:在 body 中查找 w:drawing 元素(原代码错误地用 shape_type 比较)
        has_image = False
        try:
            from docx.oxml.ns import qn

            # w:drawing 是 drawingML 图片的容器
            has_image = bool(doc.element.body.findall(f".//{qn('w:drawing')}"))
        except Exception:
            pass
        layout_type = "structured" if (doc.tables or has_header) else "plain"
        return self._success(
            {
                "layout_type": layout_type,
                "has_table": len(doc.tables) > 0,
                "has_image": has_image,
                "has_header": has_header,
                "sections": sum(1 for p in doc.paragraphs if p.style and "Heading" in p.style.name),
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            }
        )

    # ------------------------------------------------------------------
    # 内部解析器:xlsx
    # ------------------------------------------------------------------

    def _parse_xlsx(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        try:
            self._require_lib("openpyxl")
            from openpyxl import load_workbook
        except ExpertError as e:
            return self._failure(str(e))
        wb = None
        try:
            wb = load_workbook(path, read_only=True, data_only=True)
            sheets = []
            for sname in wb.sheetnames:
                ws = wb[sname]
                rows = list(ws.iter_rows(values_only=True))
                sheets.append(
                    {
                        "name": sname,
                        "rows": rows,
                        "row_count": len(rows),
                    }
                )
            metadata = {"sheet_count": len(sheets)} if extract_metadata else {}
            return self._success(
                output={
                    "type": "xlsx",
                    "metadata": metadata,
                    "sheets": sheets,
                },
                sheet_count=len(sheets),
            )
        finally:
            if wb is not None:
                try:
                    wb.close()
                except Exception:
                    pass

    def _parse_xlsx_table(self, path: str, sheet_name: str | None) -> ExpertResult:
        try:
            self._require_lib("openpyxl")
            from openpyxl import load_workbook
        except ExpertError as e:
            return self._failure(str(e))
        wb = None
        try:
            wb = load_workbook(path, read_only=True, data_only=True)
            ws = wb[sheet_name] if sheet_name else wb[wb.sheetnames[0]]
            rows = list(ws.iter_rows(values_only=True))
            headers = list(rows[0]) if rows else []
            return self._success(
                output={"headers": headers, "rows": rows[1:] if len(rows) > 1 else []},
                rows=len(rows),
            )
        except KeyError as e:
            return self._failure(f"sheet not found: {e}")
        finally:
            if wb is not None:
                try:
                    wb.close()
                except Exception:
                    pass

    def _detect_xlsx_layout(self, path: str) -> ExpertResult:
        try:
            self._require_lib("openpyxl")
            from openpyxl import load_workbook
        except ExpertError as e:
            return self._failure(str(e))
        wb = None
        try:
            wb = load_workbook(path, read_only=True, data_only=True)
            ws = wb[wb.sheetnames[0]]
            has_header = False
            row_count = 0
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                row_count += 1
                if i == 0 and all(
                    v is not None and not str(v).isdigit() for v in row if v is not None
                ):
                    has_header = True
            return self._success(
                {
                    "layout_type": "spreadsheet",
                    "has_table": row_count > 0,
                    "has_image": False,
                    "has_header": has_header,
                    "sections": 1,
                    "paragraph_count": 0,
                    "table_count": 1 if row_count > 0 else 0,
                }
            )
        finally:
            if wb is not None:
                try:
                    wb.close()
                except Exception:
                    pass

    # ------------------------------------------------------------------
    # 内部解析器:pdf
    # ------------------------------------------------------------------

    def _parse_pdf(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        # 简单实现:用 pypdf 抽文本,不抽表格(需要 pdfplumber)
        try:
            self._require_lib("pypdf")
            import pypdf
        except ExpertError as e:
            return self._failure(str(e))
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            pages_text = []
            for page in reader.pages:
                pages_text.append(page.extract_text() or "")
            metadata = {}
            if extract_metadata and reader.metadata:
                md = reader.metadata
                metadata = {
                    "title": str(md.get("/Title", "")) if md.get("/Title") else "",
                    "author": str(md.get("/Author", "")) if md.get("/Author") else "",
                    "pages": len(reader.pages),
                }
        paragraphs = []
        for pt in pages_text:
            for para in pt.split("\n\n"):
                if para.strip():
                    paragraphs.append({"text": para.strip(), "style": "Normal"})
        return self._success(
            output={
                "type": "pdf",
                "metadata": metadata,
                "paragraphs": paragraphs,
                "tables": [],
                "raw_text": "\n".join(pages_text),
            },
            paragraph_count=len(paragraphs),
        )

    def _parse_pdf_table(self, path: str, table_index: int) -> ExpertResult:
        try:
            self._require_lib("pdfplumber")
            import pdfplumber
        except ExpertError as e:
            return self._failure(
                str(e) + " (table extraction from PDF requires pdfplumber)",
            )
        with pdfplumber.open(path) as pdf:
            tables_found = 0
            for page in pdf.pages:
                tables = page.extract_tables() or []
                for tbl in tables:
                    if tables_found == table_index:
                        rows = [[(c or "") for c in r] for r in tbl]
                        headers = rows[0] if rows else []
                        return self._success(
                            output={"headers": headers, "rows": rows[1:] if len(rows) > 1 else []},
                            rows=len(rows),
                        )
                    tables_found += 1
            return self._failure(f"table_index out of range: {table_index}")

    def _detect_pdf_layout(self, path: str) -> ExpertResult:
        try:
            self._require_lib("pypdf")
            import pypdf
        except ExpertError as e:
            return self._failure(str(e))
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            page_count = len(reader.pages)
            sample_text = reader.pages[0].extract_text() if page_count else ""
        has_image = page_count > 0  # 简化:多页 PDF 通常含图
        return self._success(
            {
                "layout_type": "document",
                "has_table": False,  # 需要 pdfplumber 才能准确检测
                "has_image": has_image,
                "has_header": False,
                "sections": page_count,
                "paragraph_count": sample_text.count("\n\n") + 1 if sample_text else 0,
                "table_count": 0,
            }
        )

    # ------------------------------------------------------------------
    # 内部解析器:html
    # ------------------------------------------------------------------

    def _parse_html(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        try:
            self._require_lib("bs4")
            from bs4 import BeautifulSoup
        except ExpertError as e:
            return self._failure(str(e))
        with open(path, encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
        paragraphs = []
        for tag in soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "li"]):
            text = tag.get_text(strip=True)
            if text:
                paragraphs.append({"text": text, "style": tag.name})
        tables = []
        if extract_tables:
            for tbl in soup.find_all("table"):
                rows = []
                for tr in tbl.find_all("tr"):
                    rows.append([td.get_text(strip=True) for td in tr.find_all(["td", "th"])])
                tables.append(rows)
        metadata = {}
        if extract_metadata:
            title_tag = soup.find("title")
            metadata = {"title": title_tag.text if title_tag else ""}
        return self._success(
            output={
                "type": "html",
                "metadata": metadata,
                "paragraphs": paragraphs,
                "tables": tables,
                "raw_text": soup.get_text(separator="\n"),
            },
            paragraph_count=len(paragraphs),
            table_count=len(tables),
        )

    def _detect_html_layout(self, path: str) -> ExpertResult:
        try:
            self._require_lib("bs4")
            from bs4 import BeautifulSoup
        except ExpertError as e:
            return self._failure(str(e))
        with open(path, encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
        return self._success(
            {
                "layout_type": "web",
                "has_table": bool(soup.find("table")),
                "has_image": bool(soup.find("img")),
                "has_header": bool(soup.find(["h1", "h2", "h3"])),
                "sections": len(soup.find_all(["h1", "h2", "h3"])),
                "paragraph_count": len(soup.find_all("p")),
                "table_count": len(soup.find_all("table")),
            }
        )

    # ------------------------------------------------------------------
    # 内部解析器:pptx
    # ------------------------------------------------------------------

    def _parse_pptx(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        try:
            self._require_lib("pptx")
            from pptx import Presentation
        except ExpertError as e:
            return self._failure(str(e))
        prs = Presentation(path)
        paragraphs = []
        tables = []
        raw_parts: list[str] = []
        for slide in prs.slides:
            # 标题 placeholder(用 shape_id 判等,is 比较在 python-pptx 上失效)
            title_shape = None
            title_shape_id = None
            try:
                title_shape = slide.shapes.title
                if title_shape is not None:
                    title_shape_id = title_shape.shape_id
            except Exception:
                title_shape = None
            title = ""
            if title_shape is not None:
                try:
                    title = (title_shape.text or "").strip()
                except Exception:
                    title = ""
            if title:
                paragraphs.append({"text": title, "style": "Title"})
                raw_parts.append(title)
            for shape in slide.shapes:
                if (
                    title_shape_id is not None
                    and getattr(shape, "shape_id", None) == title_shape_id
                ):
                    continue
                if getattr(shape, "has_table", False) and extract_tables:
                    try:
                        tbl = shape.table
                        rows = [[cell.text for cell in row.cells] for row in tbl.rows]
                        tables.append(rows)
                        for row in rows:
                            raw_parts.append("\t".join(row))
                    except Exception:
                        pass
                    continue
                if shape.has_text_frame:
                    txt = shape.text_frame.text
                    if txt and txt.strip():
                        paragraphs.append({"text": txt, "style": "Normal"})
                        raw_parts.append(txt)
        metadata: dict[str, Any] = {}
        if extract_metadata:
            try:
                cp = prs.core_properties
                metadata = {
                    "title": cp.title or "",
                    "author": cp.author or "",
                    "subject": cp.subject or "",
                    "created": str(cp.created) if cp.created else "",
                    "modified": str(cp.modified) if cp.modified else "",
                    "slide_count": len(prs.slides),
                }
            except Exception:
                metadata = {"slide_count": len(prs.slides)}
        raw_text = "\n".join(raw_parts)
        return self._success(
            output={
                "type": "pptx",
                "metadata": metadata,
                "paragraphs": paragraphs,
                "tables": tables,
                "raw_text": raw_text,
                "slide_count": len(prs.slides),
            },
            paragraph_count=len(paragraphs),
            table_count=len(tables),
            slide_count=len(prs.slides),
        )

    def _detect_pptx_layout(self, path: str) -> ExpertResult:
        try:
            self._require_lib("pptx")
            from pptx import Presentation
        except ExpertError as e:
            return self._failure(str(e))
        prs = Presentation(path)
        slide_count = len(prs.slides)
        has_table = False
        has_image = False
        has_header = False
        paragraph_count = 0
        table_count = 0
        for slide in prs.slides:
            # 标题视为 header
            try:
                if slide.shapes.title is not None:
                    if (slide.shapes.title.text or "").strip():
                        has_header = True
            except Exception:
                pass
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    has_table = True
                    table_count += 1
                try:
                    if shape.shape_type == 13:
                        has_image = True
                except Exception:
                    pass
                if shape.has_text_frame:
                    if (shape.text_frame.text or "").strip():
                        paragraph_count += 1
        layout_type = "slides" if slide_count > 0 else "empty"
        return self._success(
            {
                "layout_type": layout_type,
                "has_table": has_table,
                "has_image": has_image,
                "has_header": has_header,
                "sections": slide_count,
                "paragraph_count": paragraph_count,
                "table_count": table_count,
            }
        )

    # ------------------------------------------------------------------
    # 简单格式
    # ------------------------------------------------------------------

    def _parse_text(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        with open(path, encoding="utf-8", errors="ignore") as f:
            text = f.read()
        paragraphs = [{"text": p, "style": "Normal"} for p in text.split("\n\n") if p.strip()]
        return self._success(
            output={
                "type": "text",
                "metadata": {},
                "paragraphs": paragraphs,
                "tables": [],
                "raw_text": text,
            },
            paragraph_count=len(paragraphs),
        )

    def _parse_csv(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        import csv

        with open(path, encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            rows = list(reader)
        headers = rows[0] if rows else []
        return self._success(
            output={
                "type": "csv",
                "metadata": {},
                "headers": headers,
                "rows": rows[1:] if len(rows) > 1 else [],
            },
            row_count=len(rows),
        )

    def _parse_csv_table(self, path: str) -> ExpertResult:
        import csv

        with open(path, encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            rows = list(reader)
        headers = rows[0] if rows else []
        return self._success(
            output={"headers": headers, "rows": rows[1:] if len(rows) > 1 else []},
            rows=len(rows),
        )

    def _parse_json(self, path: str, extract_tables: bool, extract_metadata: bool) -> ExpertResult:
        import json

        with open(path, encoding="utf-8") as f:
            data = json.load(f)
        return self._success(
            output={
                "type": "json",
                "metadata": {},
                "data": data,
            },
        )

    # ------------------------------------------------------------------
    # 纯文本抽取(用于表单解析)
    # ------------------------------------------------------------------

    def _extract_plain_text(self, path: str) -> ExpertResult:
        """抽取纯文本(供表单解析用)。

        Args:
            path: 文件路径

        Returns:
            ExpertResult(output=plain_text)
        """
        ext = os.path.splitext(path)[1].lstrip(".").lower()
        if ext in ("txt", "md"):
            try:
                with open(path, encoding="utf-8", errors="ignore") as f:
                    return self._success(f.read())
            except (OSError, PermissionError) as e:
                return self._failure(f"read text failed: {e}")
        # 其他格式走 parse() 取 raw_text
        r = self.parse(path, extract_tables=False, extract_metadata=False)
        if r.success:
            return self._success(r.output.get("raw_text", ""))
        return r

    # ------------------------------------------------------------------
    # Element 级别解析器
    # ------------------------------------------------------------------

    @apply_metadata
    @add_chunking_strategy
    def _parse_docx_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 docx 为 Element 列表。

        元素映射:
          - Heading 样式 → Title
          - List 样式 → ListItem
          - 普通段落 → NarrativeText
          - docx 表格 → Table
          - w:drawing → Image(占位,include_images=True 时包含)

        Args:
            path: 文件路径
            options: 解析选项

        Returns:
            Element 列表(已应用装饰器栈后处理)
        """
        try:
            self._require_lib("docx")
            from docx import Document
            from docx.oxml.ns import qn
        except ExpertError as e:
            raise OSError(f"dependency missing: {e}") from e

        doc = Document(path)
        elements: list[Element] = []
        # 遍历 body 子元素,按 XML 顺序保留段落与表格的相对位置
        body = doc.element.body
        para_idx = 0
        table_idx = 0
        paragraphs = doc.paragraphs
        tables = doc.tables
        for child in body.iterchildren():
            tag = child.tag
            if tag == qn("w:p"):
                if para_idx >= len(paragraphs):
                    continue
                p = paragraphs[para_idx]
                para_idx += 1
                text = p.text.strip()
                if not text:
                    continue
                style_name = p.style.name if p.style else "Normal"
                if "Heading" in style_name or style_name.startswith("Title"):
                    elements.append(
                        Title(
                            text=text,
                            metadata={"style": style_name, "paragraph_index": para_idx - 1},
                        )
                    )
                elif "List" in style_name:
                    elements.append(
                        ListItem(
                            text=text,
                            metadata={"style": style_name, "paragraph_index": para_idx - 1},
                        )
                    )
                else:
                    elements.append(
                        NarrativeText(
                            text=text,
                            metadata={"style": style_name, "paragraph_index": para_idx - 1},
                        )
                    )
            elif tag == qn("w:tbl") and options.extract_tables:
                if table_idx >= len(tables):
                    continue
                tbl = tables[table_idx]
                table_idx += 1
                rows = [[cell.text for cell in row.cells] for row in tbl.rows]
                elements.append(
                    Table(
                        rows=rows,
                        metadata={"table_index": table_idx - 1},
                    )
                )
            elif tag == qn("w:drawing") and options.include_images:
                elements.append(
                    Image(
                        text="",
                        metadata={"paragraph_index": para_idx - 1},
                    )
                )
        # 元数据
        if options.extract_metadata:
            cp = doc.core_properties
            file_meta = {
                "title": cp.title or "",
                "author": cp.author or "",
                "subject": cp.subject or "",
                "created": str(cp.created) if cp.created else "",
                "modified": str(cp.modified) if cp.modified else "",
            }
            for el in elements:
                el.metadata.update(file_meta)
        return elements

    @apply_metadata
    @add_chunking_strategy
    def _parse_xlsx_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 xlsx 为 Element 列表(每个 sheet 一个 Table 元素)。"""
        try:
            self._require_lib("openpyxl")
            from openpyxl import load_workbook
        except ExpertError as e:
            raise OSError(f"dependency missing: {e}") from e
        wb = None
        try:
            wb = load_workbook(path, read_only=True, data_only=True)
            elements: list[Element] = []
            for sheet_idx, sname in enumerate(wb.sheetnames):
                ws = wb[sname]
                rows = [
                    [("" if v is None else str(v)) for v in row]
                    for row in ws.iter_rows(values_only=True)
                ]
                if not rows:
                    continue
                elements.append(
                    Table(
                        rows=rows,
                        metadata={
                            "sheet_name": sname,
                            "sheet_index": sheet_idx,
                        },
                    )
                )
            return elements
        finally:
            if wb is not None:
                try:
                    wb.close()
                except Exception:
                    pass

    @apply_metadata
    @add_chunking_strategy
    def _parse_pdf_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 pdf 为 Element 列表(按页码标注 metadata.page_number)。"""
        try:
            self._require_lib("pypdf")
            import pypdf
        except ExpertError as e:
            raise OSError(f"dependency missing: {e}") from e
        elements: list[Element] = []
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page_idx, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                # 按双换行切分段落
                for para in page_text.split("\n\n"):
                    para = para.strip()
                    if not para:
                        continue
                    elements.append(
                        NarrativeText(
                            text=para,
                            metadata={"page_number": page_idx + 1},
                        )
                    )
                # 分页符标记(include_images=False 时也保留,作为页边界)
                elements.append(PageBreak(metadata={"page_number": page_idx + 1}))
            # 元数据
            if options.extract_metadata and reader.metadata:
                md = reader.metadata
                file_meta = {
                    "title": str(md.get("/Title", "")) if md.get("/Title") else "",
                    "author": str(md.get("/Author", "")) if md.get("/Author") else "",
                    "pages": len(reader.pages),
                }
                for el in elements:
                    el.metadata.update(file_meta)
        return elements

    @apply_metadata
    @add_chunking_strategy
    def _parse_html_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 html 为 Element 列表。"""
        try:
            self._require_lib("bs4")
            from bs4 import BeautifulSoup
        except ExpertError as e:
            raise OSError(f"dependency missing: {e}") from e
        with open(path, encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
        elements: list[Element] = []
        # 标题
        for tag_name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            for tag in soup.find_all(tag_name):
                text = tag.get_text(strip=True)
                if text:
                    elements.append(
                        Title(
                            text=text,
                            metadata={"tag": tag_name},
                        )
                    )
        # 列表项
        for tag in soup.find_all("li"):
            text = tag.get_text(strip=True)
            if text:
                elements.append(ListItem(text=text, metadata={"tag": "li"}))
        # 段落
        for tag in soup.find_all("p"):
            text = tag.get_text(strip=True)
            if text:
                elements.append(NarrativeText(text=text, metadata={"tag": "p"}))
        # 表格
        if options.extract_tables:
            for tbl_idx, tbl in enumerate(soup.find_all("table")):
                rows = []
                for tr in tbl.find_all("tr"):
                    rows.append([td.get_text(strip=True) for td in tr.find_all(["td", "th"])])
                if rows:
                    elements.append(
                        Table(
                            rows=rows,
                            metadata={"table_index": tbl_idx},
                        )
                    )
        # 图片
        if options.include_images:
            for img in soup.find_all("img"):
                src = img.get("src", "")
                elements.append(Image(text="", metadata={"src": src}))
        return elements

    @apply_metadata
    @add_chunking_strategy
    def _parse_text_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 txt/md 为 Element 列表(按双换行切段)。"""
        with open(path, encoding="utf-8", errors="ignore") as f:
            text = f.read()
        elements: list[Element] = []
        # Markdown 简单识别: # 开头为 Title,- * 开头为 ListItem
        for line in text.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("#"):
                # Markdown 标题
                level = len(stripped) - len(stripped.lstrip("#"))
                elements.append(
                    Title(
                        text=stripped.lstrip("# ").strip(),
                        metadata={"markdown_heading_level": level},
                    )
                )
            elif stripped.startswith(("- ", "* ", "+ ")):
                elements.append(
                    ListItem(
                        text=stripped[2:].strip(),
                        metadata={"markdown_list": True},
                    )
                )
        # 双换行切段(作为 NarrativeText 补充)
        for para in text.split("\n\n"):
            para = para.strip()
            if not para or para.startswith("#") or para.startswith(("- ", "* ", "+ ")):
                continue
            elements.append(NarrativeText(text=para))
        return elements

    @apply_metadata
    @add_chunking_strategy
    def _parse_csv_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 csv 为单个 Table 元素。"""
        import csv

        with open(path, encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            rows = [list(r) for r in reader]
        if not rows:
            return []
        return [Table(rows=rows, metadata={"row_count": len(rows)})]

    @apply_metadata
    @add_chunking_strategy
    def _parse_json_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 json 为 NarrativeText(序列化文本)。"""
        import json

        with open(path, encoding="utf-8") as f:
            data = json.load(f)
        return [
            NarrativeText(
                text=json.dumps(data, ensure_ascii=False, indent=2),
                metadata={"json_type": type(data).__name__},
            )
        ]

    @apply_metadata
    @add_chunking_strategy
    def _parse_pptx_elements(self, path: str, options: ParseOptions) -> list[Element]:
        """解析 pptx 为 Element 列表(保留幻灯片顺序)。

        元素映射:
          - 幻灯片标题 placeholder → Title
          - 文本框 → NarrativeText
          - 表格 → Table
          - 每页末尾 → PageBreak(幻灯片边界)

        metadata 含 slide_index/slide_count。
        """
        try:
            self._require_lib("pptx")
            from pptx import Presentation
        except ExpertError as e:
            raise OSError(f"dependency missing: {e}") from e
        prs = Presentation(path)
        total = len(prs.slides)
        elements: list[Element] = []
        file_meta: dict[str, Any] = {}
        if options.extract_metadata:
            try:
                cp = prs.core_properties
                file_meta = {
                    "title": cp.title or "",
                    "author": cp.author or "",
                    "subject": cp.subject or "",
                    "created": str(cp.created) if cp.created else "",
                    "modified": str(cp.modified) if cp.modified else "",
                }
            except Exception:
                pass
        for slide_idx, slide in enumerate(prs.slides):
            slide_meta = {"slide_index": slide_idx + 1, "slide_count": total}
            slide_meta.update(file_meta)
            # 标题(用 shape_id 判等,is 比较在 python-pptx 上失效)
            title_shape = None
            title_shape_id = None
            try:
                title_shape = slide.shapes.title
                if title_shape is not None:
                    title_shape_id = title_shape.shape_id
            except Exception:
                title_shape = None
            if title_shape is not None:
                try:
                    title_text = (title_shape.text or "").strip()
                except Exception:
                    title_text = ""
                if title_text:
                    elements.append(
                        Title(
                            text=title_text,
                            metadata=dict(slide_meta),
                        )
                    )
            # 其余 shapes(保留幻灯片内顺序)
            for shape in slide.shapes:
                if (
                    title_shape_id is not None
                    and getattr(shape, "shape_id", None) == title_shape_id
                ):
                    continue
                if getattr(shape, "has_table", False) and options.extract_tables:
                    try:
                        tbl = shape.table
                        rows = [[cell.text for cell in row.cells] for row in tbl.rows]
                        elements.append(
                            Table(
                                rows=rows,
                                metadata=dict(slide_meta),
                            )
                        )
                    except Exception:
                        pass
                    continue
                if shape.has_text_frame:
                    txt = shape.text_frame.text
                    if txt and txt.strip():
                        elements.append(
                            NarrativeText(
                                text=txt.strip(),
                                metadata=dict(slide_meta),
                            )
                        )
            # 幻灯片边界(与 _parse_pdf_elements 一致,每页末尾插入)
            elements.append(PageBreak(metadata=dict(slide_meta)))
        return elements
