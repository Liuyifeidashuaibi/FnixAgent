"""Template Manager(P2-9)。

文档模板管理:注册/列举/预览/应用模板,把模板变量替换为实际值。
支持 .docx / .xlsx / .pptx / .txt / .html / .md / .csv。

模板语法(统一):{{ variable_name }}
  - 支持 default:{{ name|default=Guest }}
  - 支持 conditional:{{#if condition}}...{{/if}}
  - 支持 loop:{{#each items}}...{{/each}}(简化版)

专家职责:
  - 模板注册(记录路径/扩展名/变量列表/描述)
  - 预览(用示例值填充变量返回文本)
  - 应用(渲染到 .docx/.xlsx/.pptx/.txt/.html/.md/.csv)

底层依赖:
  - python-docx / openpyxl / python-pptx(按需,可选)

降级策略:
  - 依赖缺失 → ExpertError 提示安装
  - HTML 模板渲染时对变量值做 HTML 转义,防注入
  - 输出路径扩展名校验,避免误覆盖
  - 变量替换未命中时填空串,不抛错
"""

from __future__ import annotations

import html as html_lib
import os
import re
import time
from dataclasses import dataclass, field
from typing import Any

from fnixagent.office.base import BaseExpert, ExpertError, ExpertResult

# ---------------------------------------------------------------------------
# 模板元数据
# ---------------------------------------------------------------------------


@dataclass
class TemplateInfo:
    """模板信息。

    Attributes:
        name: 模板唯一名
        path: 模板文件绝对路径
        file_type: 扩展名(docx/xlsx/pptx/txt/html/md/csv)
        description: 模板描述
        variables: 模板中识别到的变量名列表
        registered_at: 注册时间戳
    """

    name: str
    path: str
    file_type: str  # docx/xlsx/pptx/txt/html/md/csv
    description: str = ""
    variables: list[str] = field(default_factory=list)
    registered_at: float = 0.0


# ---------------------------------------------------------------------------
# TemplateManager
# ---------------------------------------------------------------------------


class TemplateManager(BaseExpert):
    """文档模板管理专家。

    全部方法返回 ExpertResult。

    能力边界:
      - 模板语法仅支持变量/default/if/each,不含复杂表达式
      - 渲染保留首个 run 样式(docx),其他格式按文本替换
      - 不支持跨模板继承
    """

    # 变量语法 {{ var }} 或 {{ var|default=xxx }}
    _VAR_PATTERN = re.compile(r"\{\{\s*([\w\.\-]+)(?:\|([^}]+))?\s*\}\}")
    # 块语法 {{#if x}}...{{/if}} / {{#each x}}...{{/each}}
    _IF_PATTERN = re.compile(r"\{\{#if\s+([\w\.\-]+)\s*\}\}(.*?)\{\{/if\}\}", re.DOTALL)
    _EACH_PATTERN = re.compile(r"\{\{#each\s+([\w\.\-]+)\s*\}\}(.*?)\{\{/each\}\}", re.DOTALL)

    # 支持的模板文件类型
    _SUPPORTED_TYPES = ("docx", "xlsx", "pptx", "txt", "html", "htm", "md", "csv")

    @property
    def name(self) -> str:
        return "template"

    def __init__(self) -> None:
        self._templates: dict[str, TemplateInfo] = {}

    # ------------------------------------------------------------------
    # 注册与列举
    # ------------------------------------------------------------------

    def register(
        self,
        name: str,
        path: str,
        description: str = "",
    ) -> ExpertResult:
        """注册一个模板。

        Args:
            name: 模板名(唯一)
            path: 模板文件路径
            description: 模板描述

        Returns:
            ExpertResult(output=TemplateInfo)
        """
        # 模板名校验
        err = self._validate_string(name, "name")
        if err:
            return self._failure(err)
        # 路径校验:必须存在 + 扩展名白名单
        err = self._validate_path(
            path,
            must_exist=True,
            allowed_exts=self._SUPPORTED_TYPES,
        )
        if err:
            return self._failure(err)
        ext = os.path.splitext(path)[1].lstrip(".").lower()
        variables = self._extract_variables(path)
        info = TemplateInfo(
            name=name,
            path=os.path.abspath(path),
            file_type=ext,
            description=description,
            variables=variables,
            registered_at=time.time(),
        )
        self._templates[name] = info
        return self._success(info, variables=variables)

    def list(self) -> ExpertResult:
        """列举所有已注册模板。

        Returns:
            ExpertResult(output=[TemplateInfo, ...])
        """
        return self._success(list(self._templates.values()))

    def get(self, name: str) -> TemplateInfo | None:
        """获取模板信息(内部使用)。"""
        return self._templates.get(name)

    # ------------------------------------------------------------------
    # 预览
    # ------------------------------------------------------------------

    def preview(
        self,
        name: str,
        sample_values: dict[str, Any] | None = None,
    ) -> ExpertResult:
        """预览模板(用示例值填充变量,返回填充后文本)。

        Args:
            name: 模板名
            sample_values: 示例值;None 用变量名占位

        Returns:
            ExpertResult(output=preview_text)
        """
        info = self._templates.get(name)
        if info is None:
            return self._failure(f"template not found: {name}")
        try:
            text = self._extract_text(info.path)
            values = sample_values or {v: f"<{v}>" for v in info.variables}
            # HTML 模板预览也做转义,与 apply 行为一致
            escape = info.file_type in ("html", "htm")
            rendered = self._render_text(text, values, escape_html=escape)
            return self._success(rendered)
        except Exception as e:
            return self._failure(f"preview failed: {e}")

    # ------------------------------------------------------------------
    # 应用
    # ------------------------------------------------------------------

    def apply(
        self,
        name: str,
        values: dict[str, Any],
        output_path: str,
    ) -> ExpertResult:
        """应用模板,生成实际文档。

        Args:
            name: 模板名
            values: 变量值字典
            output_path: 输出文件路径

        Returns:
            ExpertResult(output=output_path)
        """
        info = self._templates.get(name)
        if info is None:
            return self._failure(f"template not found: {name}")
        # 输出路径校验:扩展名须与模板同类型(max_size=None 因为是写)
        err = self._validate_path(
            output_path,
            allowed_exts=self._SUPPORTED_TYPES,
            max_size=None,
        )
        if err:
            return self._failure(err)
        out_ext = os.path.splitext(output_path)[1].lstrip(".").lower()
        if out_ext != info.file_type:
            return self._failure(
                f"output type mismatch: template is .{info.file_type} but output is .{out_ext}"
            )
        if not values:
            values = {}
        try:
            # 校验必填变量
            missing = [v for v in info.variables if v not in values]
            # missing 不阻断,用空串填充(允许部分填充)
            for v in missing:
                values.setdefault(v, "")

            ft = info.file_type
            if ft == "docx":
                return self._apply_docx(info.path, output_path, values)
            elif ft == "xlsx":
                return self._apply_xlsx(info.path, output_path, values)
            elif ft == "pptx":
                return self._apply_pptx(info.path, output_path, values)
            else:
                # 文本类(txt/html/htm/md/csv)
                return self._apply_text(info.path, output_path, values)
        except (OSError, PermissionError) as e:
            return self._failure(f"apply IO failed: {e}")
        except Exception as e:
            return self._failure(f"apply failed: {e}")

    # ------------------------------------------------------------------
    # 变量抽取
    # ------------------------------------------------------------------

    def _extract_variables(self, path: str) -> list[str]:
        """从模板文件中抽取变量名。"""
        try:
            text = self._extract_text(path)
        except Exception:
            return []
        variables: list[str] = []
        seen = set()
        # 普通变量
        for m in self._VAR_PATTERN.finditer(text):
            var = m.group(1)
            if var not in seen:
                variables.append(var)
                seen.add(var)
        # 块变量
        for pattern in (self._IF_PATTERN, self._EACH_PATTERN):
            for m in pattern.finditer(text):
                var = m.group(1)
                if var not in seen:
                    variables.append(var)
                    seen.add(var)
        return variables

    def _extract_text(self, path: str) -> str:
        """提取模板文件的文本(用于变量抽取与预览)。

        Args:
            path: 模板文件路径

        Returns:
            提取出的文本;依赖缺失或解析失败返回空串
        """
        ext = os.path.splitext(path)[1].lstrip(".").lower()
        if ext in ("txt", "html", "htm", "md", "csv"):
            try:
                with open(path, encoding="utf-8", errors="ignore") as f:
                    return f.read()
            except (OSError, PermissionError):
                return ""
        elif ext == "docx":
            try:
                from docx import Document
            except ImportError:
                return ""
            try:
                doc = Document(path)
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                return ""
        elif ext == "xlsx":
            try:
                from openpyxl import load_workbook
            except ImportError:
                return ""
            wb = None
            try:
                # read_only 模式:内存友好,适合大表
                wb = load_workbook(path, read_only=True, data_only=True)
                texts = []
                for sname in wb.sheetnames:
                    ws = wb[sname]
                    for row in ws.iter_rows(values_only=True):
                        texts.append(" ".join(str(v) for v in row if v is not None))
                return "\n".join(texts)
            except Exception:
                return ""
            finally:
                if wb is not None:
                    try:
                        wb.close()
                    except Exception:
                        pass
        elif ext == "pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return ""
            try:
                prs = Presentation(path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                return "\n".join(texts)
            except Exception:
                return ""
        return ""

    # ------------------------------------------------------------------
    # 渲染文本(支持变量/if/each)
    # ------------------------------------------------------------------

    def _render_text(
        self,
        text: str,
        values: dict[str, Any],
        *,
        escape_html: bool = False,
    ) -> str:
        """渲染纯文本模板。

        Args:
            text: 模板文本
            values: 变量值字典
            escape_html: 是否对变量值做 HTML 转义(HTML 模板应置 True 防注入)

        Returns:
            渲染后的文本
        """

        def _to_str(v: Any) -> str:
            s = "" if v is None else str(v)
            return html_lib.escape(s) if escape_html else s

        # 1. each 块
        def _each_repl(m: re.Match) -> str:
            var = m.group(1)
            body = m.group(2)
            items = values.get(var)
            if not items or not isinstance(items, (list, tuple)):
                return ""
            parts = []
            for idx, item in enumerate(items):
                ctx = dict(values)
                if isinstance(item, dict):
                    ctx.update(item)
                ctx["index"] = idx
                ctx["this"] = item
                parts.append(self._render_text(body, ctx, escape_html=escape_html))
            return "".join(parts)

        text = self._EACH_PATTERN.sub(_each_repl, text)

        # 2. if 块
        def _if_repl(m: re.Match) -> str:
            var = m.group(1)
            body = m.group(2)
            val = values.get(var)
            # 真值判断:非空/非零/非 None
            if val and val != "False" and val != "false" and val != 0:
                return self._render_text(body, values, escape_html=escape_html)
            return ""

        text = self._IF_PATTERN.sub(_if_repl, text)

        # 3. 普通变量
        def _var_repl(m: re.Match) -> str:
            var = m.group(1)
            modifier = m.group(2)
            val = values.get(var)
            if val is None:
                # default 修饰符
                if modifier and modifier.startswith("default="):
                    return modifier[len("default=") :]
                return ""
            return _to_str(val)

        text = self._VAR_PATTERN.sub(_var_repl, text)
        return text

    # ------------------------------------------------------------------
    # 应用到具体文件类型
    # ------------------------------------------------------------------

    def _apply_text(self, src: str, out: str, values: dict[str, Any]) -> ExpertResult:
        # HTML 模板需对变量值转义防注入
        ext = os.path.splitext(src)[1].lstrip(".").lower()
        escape = ext in ("html", "htm")
        try:
            with open(src, encoding="utf-8", errors="ignore") as f:
                text = f.read()
            rendered = self._render_text(text, values, escape_html=escape)
            with open(out, "w", encoding="utf-8") as f:
                f.write(rendered)
            return self._success(out, mode="text")
        except (OSError, PermissionError) as e:
            return self._failure(f"_apply_text IO failed: {e}")
        except Exception as e:
            return self._failure(f"_apply_text failed: {e}")

    def _apply_docx(self, src: str, out: str, values: dict[str, Any]) -> ExpertResult:
        try:
            self._require_lib("docx")
            from docx import Document
        except ExpertError as e:
            return self._failure(str(e))
        doc = None
        try:
            doc = Document(src)
            # 替换段落文本(保留首个 run 的样式)
            for para in doc.paragraphs:
                if "{{" in para.text:
                    new_text = self._render_text(para.text, values)
                    # 清空原 runs,写入新文本
                    if para.runs:
                        para.runs[0].text = new_text
                        for r in para.runs[1:]:
                            r.text = ""
                    else:
                        para.text = new_text
            # 替换表格单元格
            for tbl in doc.tables:
                for row in tbl.rows:
                    for cell in row.cells:
                        if "{{" in cell.text:
                            for para in cell.paragraphs:
                                if "{{" in para.text:
                                    new_text = self._render_text(para.text, values)
                                    if para.runs:
                                        para.runs[0].text = new_text
                                        for r in para.runs[1:]:
                                            r.text = ""
                                    else:
                                        para.text = new_text
            # 替换页眉页脚
            for section in doc.sections:
                for header_footer in (section.header, section.footer):
                    for para in header_footer.paragraphs:
                        if "{{" in para.text:
                            new_text = self._render_text(para.text, values)
                            if para.runs:
                                para.runs[0].text = new_text
                                for r in para.runs[1:]:
                                    r.text = ""
                            else:
                                para.text = new_text
            doc.save(out)
            return self._success(out, mode="docx")
        except (OSError, PermissionError) as e:
            return self._failure(f"_apply_docx IO failed: {e}")
        except Exception as e:
            return self._failure(f"_apply_docx failed: {e}")

    def _apply_xlsx(self, src: str, out: str, values: dict[str, Any]) -> ExpertResult:
        try:
            self._require_lib("openpyxl")
            from openpyxl import load_workbook
        except ExpertError as e:
            return self._failure(str(e))
        wb = None
        try:
            wb = load_workbook(src)
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str) and "{{" in cell.value:
                            cell.value = self._render_text(cell.value, values)
            wb.save(out)
            return self._success(out, mode="xlsx")
        except (OSError, PermissionError) as e:
            return self._failure(f"_apply_xlsx IO failed: {e}")
        except Exception as e:
            return self._failure(f"_apply_xlsx failed: {e}")
        finally:
            # 确保工作簿资源释放,避免文件句柄泄漏
            if wb is not None:
                try:
                    wb.close()
                except Exception:
                    pass

    def _apply_pptx(self, src: str, out: str, values: dict[str, Any]) -> ExpertResult:
        try:
            self._require_lib("pptx")
            from pptx import Presentation
        except ExpertError as e:
            return self._failure(str(e))
        try:
            prs = Presentation(src)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        full_text = "".join(run.text for run in para.runs)
                        if "{{" in full_text:
                            new_text = self._render_text(full_text, values)
                            if para.runs:
                                para.runs[0].text = new_text
                                for r in para.runs[1:]:
                                    r.text = ""
            prs.save(out)
            return self._success(out, mode="pptx")
        except (OSError, PermissionError) as e:
            return self._failure(f"_apply_pptx IO failed: {e}")
        except Exception as e:
            return self._failure(f"_apply_pptx failed: {e}")

    # ------------------------------------------------------------------
    # 删除模板
    # ------------------------------------------------------------------

    def unregister(self, name: str) -> ExpertResult:
        """注销模板(仅从注册表移除,不删除文件)。

        Args:
            name: 模板名

        Returns:
            ExpertResult(output=removed_path)
        """
        try:
            if name not in self._templates:
                return self._failure(f"template not found: {name}")
            info = self._templates.pop(name)
            return self._success(info.path, removed=name)
        except Exception as e:
            return self._failure(f"unregister failed: {e}")
